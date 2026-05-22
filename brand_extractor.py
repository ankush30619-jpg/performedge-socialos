"""
Brand Extractor — auto-builds a brand profile JSON from uploaded guideline files.

Supports: PDF, PNG, JPG, JPEG, WEBP, PPTX, DOCX, TXT, MD
Uses Gemini multimodal (vision + PDF) to extract:
  • brand name, tagline, category, location
  • colors, fonts
  • tone of voice, personality, voice rules
  • audience pain points & desires
  • signature phrases, do-nots
  • industry-appropriate story formula and CTA styles
"""
import io
import json
import re
import time
from pathlib import Path

import google.genai as genai
from google.genai import types as genai_types

# Supported file types
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
PDF_EXTS = {".pdf"}
PPT_EXTS = {".pptx", ".ppt"}
DOC_EXTS = {".docx", ".doc"}
TEXT_EXTS = {".txt", ".md"}


def _extract_pptx_text(path: Path) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    out = []
    try:
        prs = Presentation(str(path))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            out.append(f"\n--- Slide {slide_idx} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in p.runs).strip()
                        if text:
                            out.append(text)
    except Exception as e:
        out.append(f"[PPTX parse error: {e}]")
    return "\n".join(out)


def _extract_docx_text(path: Path) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        return ""
    try:
        doc = Document(str(path))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX parse error: {e}]"


def _extract_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[Text file read error: {e}]"


EXTRACTOR_SYSTEM_PROMPT = """You are a senior brand strategist. You will be given one or more files containing a brand's identity material — these may be PDFs, design boards, slide decks, logos, color swatches, mood boards, brand guidelines, or written notes.

Your job: extract and infer a comprehensive brand profile suitable for a social-media copywriting AI. Where the materials don't say something directly, INFER intelligently from visual style, industry conventions, and product category.

You MUST return ONLY a JSON object — no preamble, no markdown fences — matching this exact schema:

{
  "name": "Brand name",
  "tagline": "Main tagline (or inferred positioning line)",
  "alt_taglines": ["any alternate slogans found"],
  "category": "Industry/category (be specific, e.g., 'Home Appliances Brand', 'Migration Agency', 'D2C Skincare')",
  "location": "City, Region, Country if mentioned, else 'Not specified'",
  "products_or_services": ["list of products or services offered"],
  "colors": {
    "primary": "#HEXCODE",
    "secondary": "#HEXCODE",
    "accent": "#HEXCODE"
  },
  "fonts": {
    "headline": "Font name (if visible/named)",
    "body": "Font name (if visible/named)"
  },
  "audience": {
    "primary": "Detailed description of the primary target audience — demographics, geography, psychographics. Be specific.",
    "pain_points": ["5 specific real pain points this audience faces"],
    "desires": ["5 specific things this audience truly wants"]
  },
  "tone": {
    "personality": "3-5 personality traits, comma-separated (e.g., 'Bold, Aspirational, Witty, Trustworthy')",
    "voice": "Concise voice description (e.g., 'Conversational Hinglish for North Indian families' or 'Direct, expert, Australian English')",
    "rules": [
      "5-7 specific writing rules this brand should follow on social media",
      "Be specific to this brand — not generic advice"
    ]
  },
  "signature_phrases": ["5-8 on-brand phrases, taglines, or vocabulary patterns that ONLY this brand would naturally use"],
  "do_not_use": [
    "5-7 specific things this brand should NEVER do in social copy",
    "Be specific — generic platitudes, wrong tone, off-brand references etc."
  ],
  "story_formula": "One-line description of the brand's signature storytelling structure for social posts (e.g., 'Hook → Tension → Reveal → Proof → CTA')",
  "platforms_active": ["Instagram", "Facebook"],
  "cta_styles": [
    "5 example CTAs that fit this brand's voice"
  ],
  "compliance_notes": ["any regulatory or legal considerations — empty array if none"]
}

QUALITY RULES:
1. Be SPECIFIC. Do not write generic copy-bot output. Reference the actual brand category.
2. Infer colors from logos/swatches/imagery if not stated as hex codes.
3. Audience pain_points and desires must reflect the actual customer of this category, not generic.
4. Signature phrases should sound like ONLY this brand could have said them.
5. Tone rules should reflect the visual style and brand maturity you see.
6. If the brand is regional/cultural (e.g., Indian, Australian), bake that into the audience and voice.
7. If you genuinely cannot determine something, use an empty string or empty array — but TRY hard to infer first."""


def extract_brand_from_files(file_paths: list[str], api_key: str,
                             model_name: str = "gemini-2.5-flash",
                             log_callback=None) -> dict:
    """
    Given a list of file paths (PDFs, images, PPTX, DOCX, TXT, MD),
    return a structured brand profile dict via Gemini multimodal extraction.
    """
    if not file_paths:
        raise ValueError("No files provided.")
    if not api_key:
        raise ValueError("Gemini API key required.")

    _genai_client = genai.Client(api_key=api_key)

    media_parts = []  # uploaded file references for Gemini
    text_parts = []   # extracted text blobs

    for fp in file_paths:
        path = Path(fp)
        if not path.exists():
            if log_callback:
                log_callback(f"  ⚠ Skipping (not found): {fp}")
            continue

        ext = path.suffix.lower()
        if log_callback:
            log_callback(f"  • Processing: {path.name}")

        if ext in IMAGE_EXTS or ext in PDF_EXTS:
            # Upload directly to Gemini (vision/PDF support)
            try:
                up = _genai_client.files.upload(file=str(path), config={"display_name": path.name})
                # Wait for processing if needed
                for _ in range(30):
                    if up.state.name == "ACTIVE":
                        break
                    if up.state.name == "FAILED":
                        raise RuntimeError(f"Upload failed: {path.name}")
                    time.sleep(1)
                    up = _genai_client.files.get(name=up.name)
                media_parts.append(up)
                if log_callback:
                    log_callback(f"     ✓ Uploaded ({path.stat().st_size // 1024} KB)")
            except Exception as e:
                if log_callback:
                    log_callback(f"     ✗ Upload error: {e}")
                continue

        elif ext in PPT_EXTS:
            text = _extract_pptx_text(path)
            if text.strip():
                text_parts.append(f"=== From {path.name} (slide deck) ===\n{text}")
                if log_callback:
                    log_callback(f"     ✓ Extracted {len(text)} chars from PPTX")

        elif ext in DOC_EXTS:
            text = _extract_docx_text(path)
            if text.strip():
                text_parts.append(f"=== From {path.name} (document) ===\n{text}")
                if log_callback:
                    log_callback(f"     ✓ Extracted {len(text)} chars from DOCX")

        elif ext in TEXT_EXTS:
            text = _extract_text_file(path)
            if text.strip():
                text_parts.append(f"=== From {path.name} ===\n{text}")
                if log_callback:
                    log_callback(f"     ✓ Read {len(text)} chars from text file")
        else:
            if log_callback:
                log_callback(f"     ⚠ Unsupported file type: {ext}")

    if not media_parts and not text_parts:
        raise ValueError("No usable content extracted from the provided files.")

    # Build the Gemini call
    if log_callback:
        log_callback("  ⟳ Sending to Gemini for analysis…")

    config = genai_types.GenerateContentConfig(
        temperature=0.4,
        response_mime_type="application/json",
    )

    contents = [EXTRACTOR_SYSTEM_PROMPT]
    if text_parts:
        contents.append("\n\nTEXTUAL CONTENT FROM UPLOADED FILES:\n" + "\n\n".join(text_parts))
    contents.append("\n\nNow analyse all the materials (text above + media attached) and return the brand profile JSON.")
    contents.extend(media_parts)

    # Retry on rate limit
    last_err = None
    for attempt in range(4):
        try:
            resp = _genai_client.models.generate_content(
                model=model_name,
                contents=contents,
                config=config,
            )
            text = resp.text.strip()
            text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.MULTILINE).strip()
            data = json.loads(text)
            if log_callback:
                log_callback("  ✓ Brand profile extracted.")
            return _normalize_profile(data)
        except Exception as e:
            last_err = e
            err = str(e)
            wait = 0
            if "429" in err or "quota" in err.lower():
                m = re.search(r"retry[_ ]delay[^0-9]*(\d+)", err)
                wait = int(m.group(1)) + 2 if m else 25
                if log_callback:
                    log_callback(f"  ⏳ Rate limit. Waiting {wait}s…")
            else:
                wait = 2 ** attempt
            if attempt < 3:
                time.sleep(wait)
    raise RuntimeError(f"Extraction failed after retries: {last_err}")


def _normalize_profile(data: dict) -> dict:
    """Normalize the extracted profile to match the brand JSON schema used by core.py."""
    # Adapt the extracted schema to match the format core.py expects
    # core.py expects: tone.rules, tone.personality, tone.voice, audience.pain_points,
    # audience.desires, signature_phrases, do_not_use, story_formula, cta_styles, etc.
    # Our extractor already returns that schema, but we map products_or_services
    # to either products_current or services depending on category.

    products = data.get("products_or_services", [])
    category = data.get("category", "").lower()
    if "service" in category or "agency" in category or "migration" in category or "consult" in category:
        data["services"] = products
    else:
        data["products_current"] = products

    # Ensure required keys exist
    data.setdefault("alt_taglines", [])
    data.setdefault("platforms_active", ["Instagram", "Facebook"])
    data.setdefault("compliance_notes", [])
    data.setdefault("cta_styles", [])
    data.setdefault("signature_phrases", [])
    data.setdefault("do_not_use", [])
    data.setdefault("story_formula", "Hook → Tension → Reveal → Proof → CTA")

    if "colors" not in data:
        data["colors"] = {"primary": "#000000", "secondary": "#FFFFFF", "accent": "#888888"}
    if "fonts" not in data:
        data["fonts"] = {"headline": "Not specified", "body": "Not specified"}
    if "audience" not in data:
        data["audience"] = {"primary": "", "pain_points": [], "desires": []}
    if "tone" not in data:
        data["tone"] = {"personality": "", "voice": "", "rules": []}

    return data


def save_brand_profile(profile: dict, brands_dir: Path) -> str:
    """Save the extracted profile as a JSON file. Returns the key (filename without ext)."""
    name = profile.get("name", "new_brand").strip() or "new_brand"
    key = re.sub(r"[^a-z0-9_]+", "_", name.lower()).strip("_")
    if not key:
        key = "new_brand"
    out = brands_dir / f"{key}.json"
    # If exists, suffix with number
    i = 2
    while out.exists():
        out = brands_dir / f"{key}_{i}.json"
        i += 1
    out.write_text(json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")
    return out.stem
