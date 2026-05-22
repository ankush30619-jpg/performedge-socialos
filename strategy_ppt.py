"""
Polished branded PPTX builder for Growth Planner strategies.
Uses python-pptx to build an 11-slide consulting-grade deck.
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────── Color helpers ───────────
def hex_to_rgb(h):
    h = (h or "#000000").lstrip("#")
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def luminance(hex_color: str) -> float:
    h = (hex_color or "#000000").lstrip("#")
    if len(h) != 6:
        return 0.5
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0


# ─────────── Slide primitives ───────────
def _fill_solid(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(hex_color)
    shape.line.fill.background()


def _add_bg(slide, color_hex):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill_solid(bg, color_hex)
    return bg


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False, color="#222222",
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = hex_to_rgb(color)
    return box


def _add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _fill_solid(s, color)
    return s


def _add_image_safe(slide, path, x, y, *, w=None, h=None):
    if not path or not Path(path).exists():
        return None
    try:
        kw = {}
        if w: kw["width"] = w
        if h: kw["height"] = h
        return slide.shapes.add_picture(path, x, y, **kw)
    except Exception:
        return None


# ─────────── Theme ───────────
class Theme:
    def __init__(self, brand: dict, agency: dict):
        colors = brand.get("colors", {}) or {}
        self.primary = colors.get("primary") or "#0F172A"
        self.secondary = colors.get("secondary") or "#FFFFFF"
        self.accent = colors.get("accent") or "#FF6B4A"

        # Choose readable text on primary
        self.primary_text = "#FFFFFF" if luminance(self.primary) < 0.55 else "#0F172A"
        self.body_text = "#1E293B"
        self.muted = "#64748B"
        self.surface = "#F8FAFC"
        self.card = "#FFFFFF"
        self.divider = "#E2E8F0"

        # Logos
        self.brand_name = brand.get("name", "Brand")
        self.brand_tagline = brand.get("tagline", "")
        self.agency_name = agency.get("name", "PerformEdge")
        self.agency_logo_primary = agency.get("logo_primary")
        self.agency_logo_white = agency.get("logo_white")
        self.agency_logo_black = agency.get("logo_black")
        self.brand_logo = agency.get("brand_logo")  # specifically the brand's logo


# ─────────── Footer / Header ───────────
def _add_footer(slide, theme: Theme, page_num: int = None, total: int = None):
    # Thin divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5),
                                  SLIDE_H - Inches(0.55), SLIDE_W - Inches(1), Emu(8000))
    _fill_solid(div, theme.divider)

    date_str = datetime.now().strftime("%d %b %Y")
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45),
              Inches(6), Inches(0.3),
              f"Prepared by {theme.agency_name}  ·  {date_str}",
              size=9, color=theme.muted)

    # Agency logo small bottom-right
    logo = theme.agency_logo_primary or theme.agency_logo_black
    if logo:
        _add_image_safe(slide, logo,
                         SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.55),
                         h=Inches(0.4))

    if page_num and total:
        _add_text(slide, SLIDE_W - Inches(2.2), SLIDE_H - Inches(0.45),
                  Inches(0.5), Inches(0.3),
                  f"{page_num}/{total}", size=9, color=theme.muted,
                  align=PP_ALIGN.RIGHT)


def _add_top_band(slide, theme: Theme, section_label: str, slide_title: str):
    # Thin accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    _fill_solid(bar, theme.accent)

    _add_text(slide, Inches(0.5), Inches(0.32),
              Inches(8), Inches(0.35),
              section_label.upper(), size=10, bold=True,
              color=theme.muted, font="Calibri")
    _add_text(slide, Inches(0.5), Inches(0.6),
              SLIDE_W - Inches(1), Inches(1.0),
              slide_title, size=32, bold=True,
              color=theme.body_text, font="Calibri")


# ─────────── Individual slide builders ───────────
def _slide_cover(prs, theme: Theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, theme.primary)

    # Big accent stripe on left
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, Inches(0.4), SLIDE_H)
    _fill_solid(stripe, theme.accent)

    # Centered brand block
    if theme.brand_logo and Path(theme.brand_logo).exists():
        _add_image_safe(s, theme.brand_logo,
                         SLIDE_W / 2 - Inches(2), Inches(1.4),
                         w=Inches(4))

    _add_text(s, Inches(1), Inches(3.6),
              SLIDE_W - Inches(2), Inches(0.5),
              theme.brand_name.upper(),
              size=14, bold=True, color=theme.accent,
              align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(4.1),
              SLIDE_W - Inches(2), Inches(1.1),
              "Social Media\nGrowth Strategy",
              size=52, bold=True, color=theme.primary_text,
              align=PP_ALIGN.CENTER)
    if theme.brand_tagline:
        _add_text(s, Inches(1), Inches(5.7),
                  SLIDE_W - Inches(2), Inches(0.4),
                  theme.brand_tagline,
                  size=14, color="#CBD5E1",
                  align=PP_ALIGN.CENTER)

    # Bottom: prepared by
    _add_text(s, Inches(0.6), SLIDE_H - Inches(0.9),
              Inches(8), Inches(0.4),
              f"Prepared by {theme.agency_name}",
              size=10, color="#CBD5E1")
    _add_text(s, Inches(0.6), SLIDE_H - Inches(0.6),
              Inches(8), Inches(0.4),
              datetime.now().strftime("%B %Y"),
              size=10, color="#CBD5E1")

    # Agency logo bottom-right (use white variant on dark bg)
    logo = theme.agency_logo_white or theme.agency_logo_primary
    if logo:
        _add_image_safe(s, logo,
                         SLIDE_W - Inches(2), SLIDE_H - Inches(0.95),
                         h=Inches(0.55))


def _slide_section_with_columns(prs, theme: Theme,
                                  label: str, title: str, columns: list[dict],
                                  page_num, total):
    """Generic 2-3 column body layout."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, label, title)

    n = len(columns)
    if n == 0:
        _add_footer(s, theme, page_num, total)
        return s

    margin = Inches(0.5)
    gap = Inches(0.2)
    total_w = SLIDE_W - 2 * margin - (n - 1) * gap
    col_w = Emu(int(total_w / n))
    top = Inches(2.0)
    col_h = SLIDE_H - top - Inches(0.8)

    for i, c in enumerate(columns):
        x = margin + i * (col_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _fill_solid(card, theme.surface)
        card.line.color.rgb = hex_to_rgb(theme.divider)
        card.line.width = Pt(0.5)

        # Color top stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, top, col_w, Inches(0.18))
        _fill_solid(stripe, c.get("color") or theme.accent)

        _add_text(s, x + Inches(0.25), top + Inches(0.3),
                  col_w - Inches(0.5), Inches(0.6),
                  c.get("heading", ""), size=18, bold=True,
                  color=theme.body_text)
        _add_text(s, x + Inches(0.25), top + Inches(0.95),
                  col_w - Inches(0.5), col_h - Inches(1.2),
                  c.get("body", ""), size=11,
                  color=theme.body_text)

    _add_footer(s, theme, page_num, total)
    return s


def _bulleted_text(items, prefix="• "):
    if not items:
        return ""
    return "\n".join(f"{prefix}{i}" for i in items)


def _slide_executive(prs, theme: Theme, exec_data: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "01  ·  Executive Summary", "Where we are. Where we're going.")

    y = Inches(2.0)
    margin = Inches(0.5)
    half = (SLIDE_W - 2 * margin - Inches(0.3)) / 2

    # Left card: current state + opportunity
    card1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 margin, y, Emu(int(half)), Inches(3.6))
    _fill_solid(card1, theme.surface)
    card1.line.fill.background()
    _add_text(s, margin + Inches(0.3), y + Inches(0.25),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "CURRENT STATE", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), y + Inches(0.7),
              Emu(int(half - Inches(0.6))), Inches(1.4),
              exec_data.get("current_state", ""), size=13, color=theme.body_text)
    _add_text(s, margin + Inches(0.3), y + Inches(2.05),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "OPPORTUNITY", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), y + Inches(2.5),
              Emu(int(half - Inches(0.6))), Inches(1.1),
              exec_data.get("opportunity", ""), size=13, color=theme.body_text)

    # Right: big north-star + target
    x2 = margin + Emu(int(half + Inches(0.3)))
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x2, y, Emu(int(half)), Inches(3.6))
    _fill_solid(card2, theme.primary)
    card2.line.fill.background()
    _add_text(s, x2 + Inches(0.3), y + Inches(0.25),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "NORTH STAR METRIC", size=10, bold=True, color=theme.accent)
    _add_text(s, x2 + Inches(0.3), y + Inches(0.7),
              Emu(int(half - Inches(0.6))), Inches(1.2),
              exec_data.get("north_star_metric", ""), size=20, bold=True,
              color=theme.primary_text)
    _add_text(s, x2 + Inches(0.3), y + Inches(2.1),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "90-DAY TARGET", size=10, bold=True, color=theme.accent)
    _add_text(s, x2 + Inches(0.3), y + Inches(2.55),
              Emu(int(half - Inches(0.6))), Inches(1),
              exec_data.get("headline_target", ""), size=18, bold=True,
              color=theme.primary_text)

    _add_footer(s, theme, page, total)


def _slide_business(prs, theme: Theme, biz: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "02  ·  Business Understanding", "Who they are, what they're really selling.")

    y = Inches(2.0)
    margin = Inches(0.5)
    w = SLIDE_W - 2 * margin

    _add_text(s, margin, y, w, Inches(0.4),
              "WHAT THE BRAND DOES", size=10, bold=True, color=theme.accent)
    _add_text(s, margin, y + Inches(0.45), w, Inches(0.8),
              biz.get("what_brand_does", ""), size=14, color=theme.body_text)

    _add_text(s, margin, y + Inches(1.5), w, Inches(0.4),
              "UNIQUE VALUE PROPOSITION", size=10, bold=True, color=theme.accent)
    _add_text(s, margin, y + Inches(1.95), w, Inches(0.8),
              biz.get("unique_value_proposition", ""), size=14, color=theme.body_text)

    _add_text(s, margin, y + Inches(3.0), w, Inches(0.4),
              "PRIMARY BUSINESS GOALS (next 90 days)", size=10, bold=True, color=theme.accent)
    _add_text(s, margin, y + Inches(3.4), w, Inches(2),
              _bulleted_text(biz.get("primary_business_goals", [])),
              size=13, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_current_state(prs, theme: Theme, current: dict, page, total,
                           analyst_data: dict = None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")

    # Prefer Analyst data when present — that's the real Meta data
    if analyst_data:
        acc = analyst_data.get("account_metrics", {}) or {}
        prof = analyst_data.get("profile", {}) or {}
        current = {
            "instagram_handle": prof.get("username") or current.get("instagram_handle"),
            "ig_followers": acc.get("followers_now") or current.get("ig_followers"),
            "engagement_summary":
                f"Avg {acc.get('avg_engagement_rate_pct')}% ER · "
                f"{acc.get('avg_reach_pct_followers')}% reach · "
                f"{acc.get('posting_frequency_per_week')} posts/wk"
                if acc else current.get("engagement_summary", ""),
            "content_themes": current.get("content_themes", []),
            "what_is_working": current.get("what_is_working", []),
            "gaps": current.get("gaps", []),
            "facebook_handle": current.get("facebook_handle"),
        }

    _add_top_band(s, theme, f"{page:02d}  ·  Current Social State",
                  "Where the brand stands today — sourced from Analyst Baba.")

    margin = Inches(0.5); top = Inches(2.0)
    half = (SLIDE_W - 2 * margin - Inches(0.3)) / 2

    # Left card — handle + stats
    card1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 margin, top, Emu(int(half)), Inches(4.2))
    _fill_solid(card1, theme.primary)
    card1.line.fill.background()

    ig = current.get("instagram_handle") or "—"
    fb = current.get("facebook_handle") or "—"
    followers = current.get("ig_followers")
    follow_str = f"{followers:,}" if followers else "n/a"
    eng = current.get("engagement_summary") or "engagement data unavailable"

    _add_text(s, margin + Inches(0.3), top + Inches(0.25),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "INSTAGRAM", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), top + Inches(0.7),
              Emu(int(half - Inches(0.6))), Inches(0.6),
              f"@{ig}", size=22, bold=True, color=theme.primary_text)

    _add_text(s, margin + Inches(0.3), top + Inches(1.5),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "FOLLOWERS", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), top + Inches(1.9),
              Emu(int(half - Inches(0.6))), Inches(0.6),
              follow_str, size=28, bold=True, color=theme.primary_text)

    _add_text(s, margin + Inches(0.3), top + Inches(2.7),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "ENGAGEMENT", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), top + Inches(3.1),
              Emu(int(half - Inches(0.6))), Inches(0.6),
              eng, size=11, color=theme.primary_text)

    _add_text(s, margin + Inches(0.3), top + Inches(3.7),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              "FACEBOOK", size=10, bold=True, color=theme.accent)
    _add_text(s, margin + Inches(0.3), top + Inches(4.05),
              Emu(int(half - Inches(0.6))), Inches(0.4),
              fb, size=11, color=theme.primary_text)

    # Right column — themes / what's working / gaps
    x2 = margin + Emu(int(half + Inches(0.3)))
    _add_text(s, x2, top, Emu(int(half)), Inches(0.4),
              "CONTENT THEMES POSTED", size=10, bold=True, color=theme.accent)
    _add_text(s, x2, top + Inches(0.45), Emu(int(half)), Inches(1.2),
              _bulleted_text(current.get("content_themes", []) or [])[:600] or "(none yet)",
              size=10, color=theme.body_text)

    _add_text(s, x2, top + Inches(1.75), Emu(int(half)), Inches(0.4),
              "WHAT'S WORKING", size=10, bold=True, color="#10B981")
    _add_text(s, x2, top + Inches(2.2), Emu(int(half)), Inches(1.3),
              _bulleted_text(current.get("what_is_working", []) or [])[:600] or "(early-stage)",
              size=10, color=theme.body_text)

    _add_text(s, x2, top + Inches(3.6), Emu(int(half)), Inches(0.4),
              "GAPS / OPPORTUNITIES", size=10, bold=True, color="#DC2626")
    _add_text(s, x2, top + Inches(4.05), Emu(int(half)), Inches(1.3),
              _bulleted_text(current.get("gaps", []) or [])[:600] or "(strategy will address)",
              size=10, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_industry(prs, theme: Theme, ind: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "03  ·  Industry Landscape", "What's happening in this category right now.")

    y = Inches(2.0)
    margin = Inches(0.5)
    half = (SLIDE_W - 2 * margin - Inches(0.3)) / 2

    _add_text(s, margin, y, SLIDE_W - 2*margin, Inches(0.4),
              "CATEGORY STATE 2026", size=10, bold=True, color=theme.accent)
    _add_text(s, margin, y + Inches(0.45), SLIDE_W - 2*margin, Inches(1.2),
              ind.get("category_state_2026", ""), size=13, color=theme.body_text)

    cols_y = y + Inches(1.9)
    _add_text(s, margin, cols_y, Emu(int(half)), Inches(0.4),
              "KEY TRENDS", size=10, bold=True, color=theme.accent)
    _add_text(s, margin, cols_y + Inches(0.45), Emu(int(half)), Inches(2.5),
              _bulleted_text(ind.get("key_trends_in_category", [])),
              size=11, color=theme.body_text)

    x2 = margin + Emu(int(half + Inches(0.3)))
    _add_text(s, x2, cols_y, Emu(int(half)), Inches(0.4),
              "WHAT TOP BRANDS DO WELL", size=10, bold=True, color=theme.accent)
    _add_text(s, x2, cols_y + Inches(0.45), Emu(int(half)), Inches(1.4),
              _bulleted_text(ind.get("what_top_brands_do_well", [])),
              size=11, color=theme.body_text)

    _add_text(s, x2, cols_y + Inches(2.0), Emu(int(half)), Inches(0.4),
              "GAPS TO EXPLOIT", size=10, bold=True, color="#DC2626")
    _add_text(s, x2, cols_y + Inches(2.45), Emu(int(half)), Inches(1.2),
              _bulleted_text(ind.get("what_top_brands_miss", [])),
              size=11, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_audience(prs, theme: Theme, segments: list, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "04  ·  Audience Segments", "Who we're speaking to — and why they listen.")

    cols = []
    for i, seg in enumerate(segments[:3]):
        body = (f"{seg.get('description', '')}\n\n"
                f"PAINS\n" + _bulleted_text(seg.get("pain_points", []), "– ") +
                f"\n\nDESIRES\n" + _bulleted_text(seg.get("desires", []), "– ") +
                f"\n\nCONTENT THEY LOVE\n{seg.get('content_preference', '')}")
        cols.append({
            "heading": seg.get("name", f"Persona {i+1}"),
            "body": body,
            "color": theme.accent,
        })

    n = len(cols)
    if n == 0:
        _add_footer(s, theme, page, total); return
    margin = Inches(0.5); gap = Inches(0.2); top = Inches(2.0)
    total_w = SLIDE_W - 2 * margin - (n - 1) * gap
    col_w = Emu(int(total_w / n))
    col_h = SLIDE_H - top - Inches(0.8)

    for i, c in enumerate(cols):
        x = margin + i * (col_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _fill_solid(card, theme.surface)
        card.line.fill.background()
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.18))
        _fill_solid(stripe, c["color"])
        _add_text(s, x + Inches(0.25), top + Inches(0.3),
                  col_w - Inches(0.5), Inches(0.6),
                  c["heading"], size=16, bold=True, color=theme.body_text)
        _add_text(s, x + Inches(0.25), top + Inches(0.95),
                  col_w - Inches(0.5), col_h - Inches(1.2),
                  c["body"], size=10, color=theme.body_text)
    _add_footer(s, theme, page, total)


def _slide_pillars(prs, theme: Theme, pillars: list, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "05  ·  Growth Pillars", "The strategic themes that will drive every post.")

    n = min(len(pillars), 4)
    if n == 0:
        _add_footer(s, theme, page, total); return
    margin = Inches(0.5); gap = Inches(0.2); top = Inches(2.0)
    total_w = SLIDE_W - 2 * margin - (n - 1) * gap
    col_w = Emu(int(total_w / n))
    col_h = SLIDE_H - top - Inches(0.8)

    pillar_colors = [theme.accent, theme.primary, "#8B5CF6", "#10B981"]

    for i, p in enumerate(pillars[:n]):
        x = margin + i * (col_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _fill_solid(card, theme.surface)
        card.line.fill.background()
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.18))
        _fill_solid(stripe, pillar_colors[i % len(pillar_colors)])

        _add_text(s, x + Inches(0.25), top + Inches(0.3),
                  col_w - Inches(0.5), Inches(0.4),
                  f"PILLAR {i+1}", size=9, bold=True, color=theme.muted)
        _add_text(s, x + Inches(0.25), top + Inches(0.65),
                  col_w - Inches(0.5), Inches(0.8),
                  p.get("name", ""), size=14, bold=True, color=theme.body_text)
        _add_text(s, x + Inches(0.25), top + Inches(1.5),
                  col_w - Inches(0.5), Inches(1.1),
                  p.get("thesis", ""), size=10, color=theme.body_text)

        _add_text(s, x + Inches(0.25), top + Inches(2.7),
                  col_w - Inches(0.5), Inches(0.4),
                  f"{p.get('share_of_calendar', '?')}% OF CALENDAR",
                  size=9, bold=True, color=pillar_colors[i % len(pillar_colors)])
        _add_text(s, x + Inches(0.25), top + Inches(3.1),
                  col_w - Inches(0.5), Inches(0.4),
                  "FORMATS", size=9, bold=True, color=theme.muted)
        _add_text(s, x + Inches(0.25), top + Inches(3.45),
                  col_w - Inches(0.5), Inches(0.8),
                  ", ".join(p.get("content_formats", [])),
                  size=10, color=theme.body_text)
        _add_text(s, x + Inches(0.25), top + Inches(4.3),
                  col_w - Inches(0.5), Inches(0.3),
                  "EXAMPLE THEMES", size=9, bold=True, color=theme.muted)
        _add_text(s, x + Inches(0.25), top + Inches(4.6),
                  col_w - Inches(0.5), col_h - Inches(4.9),
                  _bulleted_text(p.get("example_content_themes", []), "– "),
                  size=9, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_content_strategy(prs, theme: Theme, cs: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "06  ·  Content Strategy", "Format mix, frequency, hooks.")

    margin = Inches(0.5); y = Inches(2.0)

    # Donut chart for format mix on the left
    fm = cs.get("format_mix") or {}
    if fm:
        chart_data = CategoryChartData()
        chart_data.categories = list(fm.keys())
        chart_data.add_series("Format Mix", list(fm.values()))
        chart = s.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            margin, y, Inches(5), Inches(4.5), chart_data
        ).chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.show_category_name = False
        dl.show_value = False
        dl.number_format = "0%"

    # Right side: text stats
    x2 = margin + Inches(5.3)
    w = SLIDE_W - x2 - margin
    _add_text(s, x2, y, w, Inches(0.4),
              "POSTING FREQUENCY", size=10, bold=True, color=theme.accent)
    _add_text(s, x2, y + Inches(0.45), w, Inches(0.7),
              cs.get("posting_frequency", ""), size=13, color=theme.body_text)

    _add_text(s, x2, y + Inches(1.4), w, Inches(0.4),
              "HOOK STRATEGY", size=10, bold=True, color=theme.accent)
    _add_text(s, x2, y + Inches(1.85), w, Inches(1.2),
              cs.get("hook_strategy", ""), size=11, color=theme.body_text)

    _add_text(s, x2, y + Inches(3.2), w, Inches(0.4),
              "HASHTAG STRATEGY", size=10, bold=True, color=theme.accent)
    _add_text(s, x2, y + Inches(3.65), w, Inches(1),
              cs.get("hashtag_strategy", ""), size=11, color=theme.body_text)

    _add_text(s, margin, SLIDE_H - Inches(1.2), SLIDE_W - 2*margin, Inches(0.5),
              f"Best posting times: {cs.get('best_posting_times', 'TBD')}",
              size=11, color=theme.muted)

    _add_footer(s, theme, page, total)


def _slide_roadmap(prs, theme: Theme, roadmap: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, f"{page:02d}  ·  30-Day Campaign Roadmap", "Week-by-week arc.")

    margin = Inches(0.5)
    top = Inches(2.0)
    n = 4
    gap = Inches(0.2)
    total_w = SLIDE_W - 2 * margin - (n - 1) * gap
    col_w = Emu(int(total_w / n))
    col_h = SLIDE_H - top - Inches(0.8)

    # Support both new (weeks) + legacy (months) schemas
    if any(k.startswith("week_") for k in (roadmap.keys() if isinstance(roadmap, dict) else [])):
        rows = [
            ("WEEK 1", roadmap.get("week_1_theme", ""), roadmap.get("week_1_campaigns", [])),
            ("WEEK 2", roadmap.get("week_2_theme", ""), roadmap.get("week_2_campaigns", [])),
            ("WEEK 3", roadmap.get("week_3_theme", ""), roadmap.get("week_3_campaigns", [])),
            ("WEEK 4", roadmap.get("week_4_theme", ""), roadmap.get("week_4_campaigns", [])),
        ]
    else:
        rows = [
            ("MONTH 1", roadmap.get("month_1_theme", ""), roadmap.get("month_1_campaigns", [])),
            ("MONTH 2", roadmap.get("month_2_theme", ""), roadmap.get("month_2_campaigns", [])),
            ("MONTH 3", roadmap.get("month_3_theme", ""), roadmap.get("month_3_campaigns", [])),
        ]
        n = 3
        total_w = SLIDE_W - 2 * margin - (n - 1) * gap
        col_w = Emu(int(total_w / n))

    for i, (label, theme_txt, campaigns) in enumerate(rows):
        x = margin + i * (col_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _fill_solid(card, theme.surface)
        card.line.fill.background()
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.18))
        _fill_solid(stripe, theme.accent)
        _add_text(s, x + Inches(0.25), top + Inches(0.3),
                  col_w - Inches(0.5), Inches(0.4),
                  label, size=9, bold=True, color=theme.muted)
        _add_text(s, x + Inches(0.25), top + Inches(0.7),
                  col_w - Inches(0.5), Inches(1.0),
                  theme_txt, size=13, bold=True, color=theme.body_text)
        _add_text(s, x + Inches(0.25), top + Inches(1.95),
                  col_w - Inches(0.5), Inches(0.4),
                  "KEY CAMPAIGNS", size=9, bold=True, color=theme.accent)
        _add_text(s, x + Inches(0.25), top + Inches(2.35),
                  col_w - Inches(0.5), col_h - Inches(2.7),
                  _bulleted_text(campaigns), size=10, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_version_changes(prs, theme: Theme, strategy: dict, page, total):
    """v2+ only — what changed from previous version."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    meta = strategy.get("_meta", {}) or {}
    v = meta.get("version", 2); prev = meta.get("previous_version", v-1)
    _add_top_band(s, theme,
                  f"{page:02d}  ·  What's New in v{v}",
                  f"Specific changes from v{prev} — driven by Analyst Baba's data.")

    changes = strategy.get("changes_from_previous_version", []) or []
    margin = Inches(0.5); y = Inches(2.0)
    w = SLIDE_W - 2 * margin
    item_h = Inches(0.6)

    if not changes:
        _add_text(s, margin, y, w, Inches(2),
                  "(No version-specific changes were recorded by the strategist.\n"
                  "The plan may have changed substantially — review side-by-side with v"
                  + str(prev) + ".)",
                  size=12, color=theme.muted)
    else:
        for i, change in enumerate(changes[:8]):
            yi = y + Inches(0.05) + i * (item_h + Inches(0.1))
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, margin, yi + Inches(0.08),
                                         Inches(0.45), Inches(0.45))
            _fill_solid(badge, theme.accent)
            _add_text(s, margin, yi + Inches(0.08), Inches(0.45), Inches(0.45),
                      str(i+1), size=12, bold=True, color="#FFFFFF",
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text(s, margin + Inches(0.6), yi, w - Inches(0.6), item_h,
                      change, size=11, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_trends(prs, theme: Theme, trends: list, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "08  ·  Trend Watch", "Specific trends to ride this month.")

    margin = Inches(0.5); y = Inches(2.0)
    w = SLIDE_W - 2 * margin
    item_h = Inches(1.0)
    for i, t in enumerate(trends[:5]):
        yi = y + Inches(0.1) + i * (item_h + Inches(0.15))
        # Numbered badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, margin, yi, Inches(0.6), Inches(0.6))
        _fill_solid(badge, theme.accent)
        _add_text(s, margin, yi, Inches(0.6), Inches(0.6),
                  str(i+1), size=18, bold=True, color="#FFFFFF",
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Text block
        x_text = margin + Inches(0.8)
        _add_text(s, x_text, yi - Inches(0.05), w - Inches(0.8), Inches(0.4),
                  t.get("trend", ""), size=13, bold=True, color=theme.body_text)
        _add_text(s, x_text, yi + Inches(0.35), w - Inches(0.8), Inches(0.35),
                  f"WHY RELEVANT: {t.get('why_relevant', '')}",
                  size=9, color=theme.muted)
        _add_text(s, x_text, yi + Inches(0.7), w - Inches(0.8), Inches(0.4),
                  f"How to ride: {t.get('how_to_ride', '')}",
                  size=10, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_kpis(prs, theme: Theme, kpis: dict, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "09  ·  KPIs & Success Metrics", "What success looks like in 90 days.")

    margin = Inches(0.5); top = Inches(2.0)
    items = [
        ("FOLLOWER GROWTH", kpis.get("follower_growth_target_90_days", "")),
        ("ENGAGEMENT RATE", kpis.get("engagement_rate_target", "")),
        ("REACH", kpis.get("reach_target", "")),
        ("SAVES + SHARES", kpis.get("save_share_target", "")),
        ("CONVERSION", kpis.get("conversion_target", "")),
    ]
    cols = 3
    rows = 2
    gap = Inches(0.2)
    total_w = SLIDE_W - 2 * margin - (cols - 1) * gap
    col_w = Emu(int(total_w / cols))
    row_h = Inches(2.1)

    for i, (label, val) in enumerate(items):
        col = i % cols
        row = i // cols
        x = margin + col * (col_w + gap)
        y = top + row * (row_h + Inches(0.2))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
        _fill_solid(card, theme.primary)
        card.line.fill.background()
        _add_text(s, x + Inches(0.25), y + Inches(0.25),
                  col_w - Inches(0.5), Inches(0.4),
                  label, size=10, bold=True, color=theme.accent)
        _add_text(s, x + Inches(0.25), y + Inches(0.7),
                  col_w - Inches(0.5), row_h - Inches(0.9),
                  val or "TBD", size=14, bold=True, color=theme.primary_text)

    _add_footer(s, theme, page, total)


def _slide_recommendations(prs, theme: Theme, recs: list, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, "#FFFFFF")
    _add_top_band(s, theme, "10  ·  Recommendations", "Specific next steps, in priority order.")

    margin = Inches(0.5); y = Inches(2.0)
    w = SLIDE_W - 2 * margin
    item_h = Inches(0.55)
    for i, r in enumerate(recs[:8]):
        yi = y + Inches(0.05) + i * (item_h + Inches(0.1))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, margin, yi + Inches(0.05),
                                     Inches(0.45), Inches(0.45))
        _fill_solid(badge, theme.accent)
        _add_text(s, margin, yi + Inches(0.05), Inches(0.45), Inches(0.45),
                  str(i+1), size=12, bold=True, color="#FFFFFF",
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(s, margin + Inches(0.6), yi, w - Inches(0.6), item_h,
                  r, size=12, color=theme.body_text)

    _add_footer(s, theme, page, total)


def _slide_closing(prs, theme: Theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, theme.primary)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    _fill_solid(stripe, theme.accent)

    _add_text(s, Inches(1), Inches(3.0),
              SLIDE_W - Inches(2), Inches(1),
              "Let's build.",
              size=72, bold=True, color=theme.primary_text,
              align=PP_ALIGN.CENTER)
    _add_text(s, Inches(1), Inches(4.5),
              SLIDE_W - Inches(2), Inches(0.5),
              f"{theme.agency_name}  ·  Social Growth Partners",
              size=14, color="#CBD5E1",
              align=PP_ALIGN.CENTER)

    logo = theme.agency_logo_white or theme.agency_logo_primary
    if logo:
        _add_image_safe(s, logo,
                         SLIDE_W / 2 - Inches(1.5), SLIDE_H - Inches(1.5),
                         w=Inches(3))


# ─────────── Public builder ───────────
def build_strategy_pptx(strategy: dict, brand: dict, settings: dict,
                          output_path: str, brand_logo: str = None) -> str:
    """Build the full polished strategy deck and save to output_path."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    agency = {
        "name": settings.get("agency_name", "PerformEdge"),
        "logo_primary": settings.get("agency_logo_primary"),
        "logo_white": settings.get("agency_logo_white"),
        "logo_black": settings.get("agency_logo_black"),
        "brand_logo": brand_logo,
    }
    theme = Theme(brand, agency)

    # Slide count adapts to mode + version
    meta = strategy.get("_meta") or {}
    mode = meta.get("brand_mode", "existing")
    version = meta.get("version", 1)
    show_current_state = (mode == "existing")
    show_version_diff = (version > 1)

    # Compute total dynamically
    base = 11  # cover + 9 content slides + closing equivalents (we count content only for footer)
    total = base + (1 if show_current_state else 0) + (1 if show_version_diff else 0)

    page = 1
    _slide_cover(prs, theme)
    page += 1
    _slide_executive(prs, theme, strategy.get("executive_summary", {}), page, total); page += 1
    if show_current_state:
        # Try to load latest Analyst report for richer data
        analyst_data = None
        try:
            import analyst_agent as _aa
            import core as _core
            brand_key = meta.get("brand_key")
            if brand_key:
                analyst_data = _aa.latest_report(_core.BRANDS_DIR, brand_key)
        except Exception:
            analyst_data = None
        _slide_current_state(prs, theme, strategy.get("current_social_state", {}),
                              page, total, analyst_data=analyst_data); page += 1
    if show_version_diff:
        _slide_version_changes(prs, theme, strategy, page, total); page += 1
    _slide_business(prs, theme, strategy.get("business_understanding", {}), page, total); page += 1
    _slide_industry(prs, theme, strategy.get("industry_landscape", {}), page, total); page += 1
    _slide_audience(prs, theme, strategy.get("audience_segments", []) or [], page, total); page += 1
    _slide_pillars(prs, theme, strategy.get("growth_pillars", []) or [], page, total); page += 1
    _slide_content_strategy(prs, theme, strategy.get("content_strategy", {}), page, total); page += 1
    _slide_roadmap(prs, theme, strategy.get("campaign_roadmap", {}), page, total); page += 1
    _slide_trends(prs, theme, strategy.get("trend_watch", []) or [], page, total); page += 1
    _slide_kpis(prs, theme, strategy.get("kpis", {}), page, total); page += 1
    _slide_recommendations(prs, theme, strategy.get("recommendations_priority_order", []) or [], page, total); page += 1
    _slide_closing(prs, theme)

    prs.save(output_path)
    return output_path
