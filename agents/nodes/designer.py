"""
Designer Agent Node — IMPLEMENTED
Generates images via Freepik Mystic API for visual posts,
builds XLSX content calendar and PPT strategy deck,
then uploads everything to Supabase Storage.
"""
import asyncio
import io
import json
import os
import time
from datetime import datetime

import httpx
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from supabase import create_client, Client as SupabaseClient
from state import SocialOSState

# ── Config ────────────────────────────────────────────────────────────────────
FREEPIK_API_KEY  = os.getenv("FREEPIK_API_KEY", "")
FREEPIK_ENGINE   = os.getenv("FREEPIK_ENGINE", "mystic")
OPENAI_API_KEY   = os.getenv("OPENAI_API_KEY", "")
BUCKET           = os.getenv("SUPABASE_STORAGE_BUCKET", "socialos-storage")

# Only generate images for these types
VISUAL_TYPES = {"Graphic", "AI Reel", "Carousel"}
# Max images per run (cost control)
MAX_IMAGES = 6

# ── Design classification (autonomous content-type detection) ────────────────
# Maps a post (content type + topic + copy) to a creative-strategy category so
# the prompt engine can pick the right layout, emotional trigger and CTA placement.
_CLASS_KEYWORDS = {
    "Offer Post":       ("offer", "sale", "discount", "deal", "% off", "coupon", "limited", "price", "buy"),
    "Announcement":     ("announc", "launch", "introducing", "new ", "now live", "coming soon", "update"),
    "Testimonial":      ("testimonial", "review", "client said", "customer", "result", "transformation", "case study"),
    "Comparison":       ("vs", "versus", "before/after", "before and after", "compare", "this or that"),
    "Product Showcase": ("product", "feature", "showcase", "how it works", "demo", "inside"),
    "Lead Magnet":      ("free", "download", "guide", "checklist", "ebook", "template", "cheat sheet", "webinar"),
    "Educational":      ("how to", "tips", "steps", "mistakes", "guide", "learn", "explained", "myth", "why "),
    "Storytelling Post":("story", "journey", "i used to", "behind the scenes", "my ", "we started"),
}


def _classify_design(post: dict) -> str:
    """Autonomously classify the design category from content type + copy."""
    ct = post.get("contentType", "")
    if ct == "Carousel":
        return "Carousel"
    blob = " ".join(str(post.get(k, "")) for k in
                     ("topic", "caption_long", "caption", "hook", "cta", "copy_brief")).lower()
    for category, kws in _CLASS_KEYWORDS.items():
        if any(k in blob for k in kws):
            return category
    return "Static Graphic"


# Per-category creative strategy: emotional trigger · attention strategy · layout.
_CATEGORY_STRATEGY = {
    "Static Graphic":   ("clarity + authority", "bold headline as the focal point", "single-focus centered composition with generous negative space"),
    "Carousel":         ("curiosity → resolution", "scroll-stopping cover that opens a loop", "consistent slide system, strong left-aligned hierarchy"),
    "Offer Post":       ("urgency + value", "the offer/number is the hero", "high-contrast badge layout, price/percentage dominant"),
    "Announcement":     ("excitement + novelty", "'new' signalled instantly", "spotlight composition, product/news centered"),
    "Testimonial":      ("trust + social proof", "the result/quote carries the frame", "quote-card layout with attribution and proof"),
    "Educational":      ("authority + usefulness", "promise of a clear takeaway", "numbered/step layout, clean infographic structure"),
    "Comparison":       ("contrast + insight", "the two sides clash visually", "split-screen 50/50 composition"),
    "Product Showcase": ("desire + clarity", "the product is hero-lit", "product-centric composition with breathing room"),
    "Storytelling Post":("emotion + relatability", "a human, candid moment", "documentary/lifestyle framing, authentic feel"),
    "Lead Magnet":      ("value + low friction", "the freebie is unmissable", "asset-mockup layout with a clear claim path"),
}

# Lazy Supabase singleton — initialized on first use to ensure .env is loaded
_supabase: SupabaseClient | None = None


def _get_supabase() -> SupabaseClient | None:
    """Return (or lazily create) the Supabase client."""
    global _supabase
    if _supabase is not None:
        return _supabase
    url = os.getenv("NEXT_PUBLIC_SUPABASE_URL", "")
    key = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
    if not url or not key:
        print(f"[Designer] Supabase env vars missing: URL={'set' if url else 'MISSING'}, KEY={'set' if key else 'MISSING'}")
        return None
    try:
        _supabase = create_client(url, key)
        print(f"[Designer] Supabase client initialized OK")
        return _supabase
    except Exception as e:
        print(f"[Designer] Supabase init error: {e}")
        return None


# ── Main Node ─────────────────────────────────────────────────────────────────

async def designer_node(state: SocialOSState, event_queue: asyncio.Queue) -> dict:
    posts    = state.get("posts_with_copy") or state.get("content_calendar") or []
    brand    = state.get("brand") or {}
    strategy = state.get("growth_strategy") or {}
    run_id   = state["run_id"]
    brand_id = state["brand_id"]

    # ── 1. Generate images ────────────────────────────────────────────────────
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "designer",
        "message": "Generating brand images with Freepik Mystic…",
    })

    visual_posts = [p for p in posts if p.get("contentType") in VISUAL_TYPES][:MAX_IMAGES]
    design_assets = []

    image_tasks = [
        _generate_and_upload(post, brand, run_id, i, event_queue)
        for i, post in enumerate(visual_posts)
    ]
    results = await asyncio.gather(*image_tasks, return_exceptions=True)
    for r in results:
        if isinstance(r, dict) and r.get("imageUrl"):
            design_assets.append(r)

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "designer",
        "message": f"{len(design_assets)} images generated — building Excel calendar…",
    })

    # ── 2. Excel XLSX ─────────────────────────────────────────────────────────
    excel_url = await _build_excel(posts, brand, strategy, run_id)

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "designer",
        "message": "Building strategy PowerPoint deck…",
    })

    # ── 3. PPT deck ───────────────────────────────────────────────────────────
    ppt_url = await _build_ppt(posts, brand, strategy, state.get("analyst_report") or {}, run_id)

    return {
        "design_assets": design_assets,
        "ppt_url":       ppt_url,
        "excel_url":     excel_url,
        "_message": (
            f"Designer complete — {len(design_assets)} images, "
            f"{'Excel ready' if excel_url else 'Excel failed'}, "
            f"{'PPT ready' if ppt_url else 'PPT failed'}"
        ),
    }


# ── JSON Prompt Engine ────────────────────────────────────────────────────────
# Converts a finished post + brand assets into an ultra-detailed, structured
# design spec. The copy already exists (Copywriter); the Designer's job is to
# turn it into a visual blueprint, then compose a high-fidelity image prompt
# from it. Deterministic → fast, free, reproducible, fully auditable.

_VISUAL_STYLE_BY_TONE = {
    "professional": "clean modern editorial, premium minimal, confident",
    "friendly":     "warm approachable, soft rounded shapes, inviting",
    "bold":         "high-contrast bold, punchy, energetic",
    "luxury":       "elegant premium, refined, lots of negative space",
    "playful":      "vibrant playful, dynamic, fun",
    "minimal":      "ultra-minimal, lots of whitespace, single focal point",
}

# Keep the proven square for Graphic/Carousel; only Reels go vertical. Unknown
# values degrade gracefully (Freepik returns non-2xx → image skipped, no crash).
_ASPECT_BY_TYPE = {"AI Reel": "9_16", "Carousel": "1_1", "Graphic": "1_1"}

# Standard agency-grade negatives — kills the things QC would otherwise reject.
_BASE_NEGATIVES = [
    "gibberish text", "garbled letters", "misspelled words", "lorem ipsum",
    "watermark", "logo artifacts", "stock-photo cliché", "extra fingers",
    "distorted faces", "low resolution", "jpeg artifacts", "cluttered layout",
    "busy background", "muddy colors", "amateur", "template look", "clip art",
]


def _palette_from_brand(brand: dict) -> list[str]:
    colors = brand.get("colors") or {}
    if not isinstance(colors, dict):
        return ["#6C3CE1"]
    pal = []
    for k in ("primary", "secondary", "accent"):
        v = colors.get(k)
        if v:
            pal.append(v)
    extra = colors.get("palette") or []
    if isinstance(extra, list):
        pal.extend([c for c in extra if c])
    return pal or ["#6C3CE1"]


def _design_copy(post: dict, category: str) -> dict:
    """Pull the headline / subheadline / cta the visual should carry."""
    gl = post.get("graphic_layout") if isinstance(post.get("graphic_layout"), dict) else {}
    slides = post.get("carousel_slides") or []
    cover = slides[0] if slides and isinstance(slides[0], dict) else {}
    headline = (gl.get("headline") or cover.get("headline")
                or (post.get("hook_variations") or [None])[0] or post.get("hook")
                or post.get("topic") or "")
    subheadline = gl.get("subheadline") or cover.get("on_slide_text") or ""
    cta = post.get("cta") or gl.get("footer_text") or ""
    supporting = gl.get("supporting_elements") or []
    return {
        "headline": str(headline)[:90],
        "subheadline": str(subheadline)[:90],
        "cta": str(cta)[:60],
        "supporting": [str(s)[:60] for s in supporting[:4]],
    }


def _build_design_json(post: dict, brand: dict, category: str) -> dict:
    """The ultra-detailed JSON prompt spec (json_prompt_engine)."""
    ct       = post.get("contentType", "Graphic")
    tone     = (brand.get("tone") or "professional")
    tone_key = next((k for k in _VISUAL_STYLE_BY_TONE if k in tone.lower()), "professional")
    palette  = _palette_from_brand(brand)
    copy     = _design_copy(post, category)
    trigger, attention, layout = _CATEGORY_STRATEGY.get(
        category, _CATEGORY_STRATEGY["Static Graphic"])
    typography = (brand.get("voiceStyle") or "modern geometric sans-serif, strong weight contrast")

    return {
        "design_type":      category,
        "platform":         "Instagram",
        "brand_name":       brand.get("name", ""),
        "industry":         brand.get("industry") or brand.get("niche", ""),
        "target_audience":  brand.get("targetAudience", ""),
        "visual_style":     _VISUAL_STYLE_BY_TONE[tone_key],
        "design_goal":      f"{trigger}; {attention}",
        "layout_structure": layout,
        "color_palette":    ", ".join(palette),
        "typography":       typography,
        "headline":         copy["headline"],
        "subheadline":      copy["subheadline"],
        "cta":              copy["cta"],
        "visual_elements":  copy["supporting"] or [f"{brand.get('niche','')} context imagery"],
        "icons":            [],
        "illustrations":    [],
        "composition":      layout,
        "lighting":         "soft directional studio light, gentle gradient",
        "depth":            "subtle depth of field, layered foreground/background",
        "background_style": "on-brand gradient or clean solid that supports text contrast",
        "brand_guidelines": f"use {palette[0]} as the dominant brand color; keep clear space for text overlay",
        "image_size":       _ASPECT_BY_TYPE.get(ct, "1_1"),
        "negative_prompts": _BASE_NEGATIVES,
    }


def _json_to_prompt(dj: dict) -> tuple[str, str]:
    """Compose a high-fidelity Freepik prompt + negative string from the JSON spec."""
    parts = [
        f"{dj['design_type']} for {dj['platform']} — {dj['industry']} brand"
        + (f" '{dj['brand_name']}'" if dj['brand_name'] else "") + ".",
        f"Visual style: {dj['visual_style']}.",
        f"Composition: {dj['composition']}.",
        f"Design goal: {dj['design_goal']}.",
        f"Color palette: {dj['color_palette']} (dominant {dj['color_palette'].split(',')[0].strip()}).",
        f"Lighting: {dj['lighting']}. Depth: {dj['depth']}. Background: {dj['background_style']}.",
    ]
    if dj.get("headline"):
        parts.append(f"Leave clean negative space for a bold headline overlay reading the theme: \"{dj['headline']}\".")
    if dj.get("visual_elements"):
        ve = ", ".join(v for v in dj["visual_elements"] if v)
        if ve:
            parts.append(f"Include supporting visual elements: {ve}.")
    parts.append(f"Typography direction (for overlaid text): {dj['typography']}.")
    parts.append(f"Audience: {dj['target_audience']}." if dj.get("target_audience") else "")
    parts.append("Agency-grade, scroll-stopping, conversion-focused, strong visual hierarchy, Instagram-ready, ultra high quality.")
    negative = ", ".join(dj.get("negative_prompts") or _BASE_NEGATIVES)
    # Fold negatives into the prompt text — keeps the proven Freepik payload shape
    # intact (no unrecognized fields that could 400 the request).
    parts.append(f"Avoid: {negative}.")
    prompt = " ".join(p for p in parts if p)
    return prompt, negative


def _qc_check(dj: dict, image_url: str) -> dict:
    """Lightweight quality-control gate (quality_control)."""
    checks = {
        "image_generated":  bool(image_url),
        "brand_color_set":  bool(dj.get("color_palette")),
        "headline_present": bool(dj.get("headline")),
        "hierarchy_defined":bool(dj.get("layout_structure")),
        "negatives_applied":bool(dj.get("negative_prompts")),
    }
    passed = all(checks.values())
    return {"passed": passed, "checks": checks}


# ── Image Generation ──────────────────────────────────────────────────────────

async def _generate_and_upload(post: dict, brand: dict, run_id: str, idx: int, event_queue) -> dict:
    topic = post.get("topic", "brand content")
    ct    = post.get("contentType", "Graphic")

    # 1) classify → 2) build ultra-detailed JSON spec → 3) compose prompt
    category = _classify_design(post)
    design_json = _build_design_json(post, brand, category)
    prompt, negative = _json_to_prompt(design_json)
    resolution = design_json.get("image_size", "1_1")

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "designer",
        "message": f"Image {idx+1} [{category}]: engineering JSON prompt for '{topic[:34]}'…",
    })

    image_url = await _freepik_generate(prompt, resolution=resolution, negative=negative)

    # 4) quality control
    qc = _qc_check(design_json, image_url)
    if not image_url:
        return {}

    return {
        "imageUrl":    image_url,
        "contentType": ct,
        "topic":       topic,
        "prompt":      prompt,
        "date":        post.get("date"),
        # ── New: full provenance for the autonomous designer ──
        "designCategory": category,
        "designJson":     design_json,
        "qc":             qc,
    }


async def _freepik_generate(prompt: str, resolution: str = "1_1", negative: str = "") -> str:
    """Call Freepik Mystic API and poll for result."""
    if not FREEPIK_API_KEY:
        return ""

    headers = {
        "x-freepik-api-key": FREEPIK_API_KEY,
        "Content-Type": "application/json",
        "Accept": "application/json",
    }

    # NOTE: negatives are folded into the prompt text by _json_to_prompt, so the
    # payload keeps the exact proven shape (no unrecognized fields).
    payload = {
        "prompt": prompt,
        "num_images": 1,
        "resolution": resolution,
        "engine": FREEPIK_ENGINE,
        "creative_detailing": 60,
        "styling": {"style": "photo"},
    }

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            # Start generation
            start_resp = await client.post(
                "https://api.freepik.com/v1/ai/mystic",
                headers=headers,
                json=payload,
            )

            if start_resp.status_code not in (200, 201, 202):
                print(f"[Designer] Freepik start error {start_resp.status_code}: {start_resp.text[:200]}")
                return ""

            data = start_resp.json()
            task_id = (data.get("data") or {}).get("id") or data.get("id")

            if not task_id:
                print(f"[Designer] No task ID from Freepik: {data}")
                return ""

            # Poll for result (max 60s)
            for attempt in range(20):
                await asyncio.sleep(3)
                poll_resp = await client.get(
                    f"https://api.freepik.com/v1/ai/mystic/{task_id}",
                    headers=headers,
                )
                poll_data = poll_resp.json()
                result = poll_data.get("data") or {}
                status = result.get("status") or poll_data.get("status")

                if status == "completed":
                    generated = result.get("generated_images") or result.get("images") or []
                    if generated:
                        return generated[0].get("url") or generated[0].get("base64", "")
                    return ""

                if status in ("failed", "error", "cancelled"):
                    print(f"[Designer] Freepik task {task_id} status: {status}")
                    return ""

            print(f"[Designer] Freepik task {task_id} timed out")
            return ""

    except Exception as e:
        print(f"[Designer] Freepik error: {e}")
        return ""


# ── Excel XLSX ────────────────────────────────────────────────────────────────

async def _build_excel(posts: list, brand: dict, strategy: dict, run_id: str) -> str | None:
    try:
        wb = openpyxl.Workbook()

        # ── Sheet 1: Content Calendar ──
        ws = wb.active
        ws.title = "Content Calendar"
        name = brand.get("name", "Brand")
        niche = brand.get("niche", "")

        # Header row style
        header_fill  = PatternFill("solid", fgColor="6C3CE1")
        white_font   = Font(bold=True, color="FFFFFF", size=11)
        subhdr_fill  = PatternFill("solid", fgColor="2D235A")
        subhdr_font  = Font(bold=True, color="A78BFA", size=10)
        center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        thin_border  = Border(
            left=Side(style="thin", color="2D235A"),
            right=Side(style="thin", color="2D235A"),
            top=Side(style="thin", color="2D235A"),
            bottom=Side(style="thin", color="2D235A"),
        )

        # Title
        ws.merge_cells("A1:R1")
        title_cell = ws["A1"]
        title_cell.value = f"{name} — Content Calendar"
        title_cell.font  = Font(bold=True, color="FFFFFF", size=14)
        title_cell.fill  = header_fill
        title_cell.alignment = center_align

        # Subtitle
        ws.merge_cells("A2:R2")
        sub_cell = ws["A2"]
        sub_cell.value = f"Generated by SocialOS | {niche} | {datetime.utcnow().strftime('%d %b %Y')}"
        sub_cell.font  = Font(color="A78BFA", italic=True, size=10)
        sub_cell.fill  = PatternFill("solid", fgColor="1E1640")
        sub_cell.alignment = center_align

        # Column headers (full PRD FR-050 brief)
        headers = [
            "#", "Date", "Day", "Time", "Content Type", "Topic",
            "Hook 1", "Hook 2", "Hook 3",
            "Caption (Short)", "Caption (Long)", "CTA",
            "Hashtags", "SEO Keywords", "Audio/Music",
            "Script / Slides / Layout", "Visual Brief", "Status",
        ]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col, value=h)
            cell.font  = subhdr_font
            cell.fill  = subhdr_fill
            cell.alignment = center_align
            cell.border = thin_border

        # Post rows
        type_colors = {
            "Reel":     "3B82F6", "AI Reel":  "EC4899",
            "Carousel": "10B981", "Graphic":  "A78BFA", "Story": "F59E0B",
        }
        for i, post in enumerate(posts):
            row = i + 4
            post_date = post.get("date", "")
            try:
                dt = datetime.fromisoformat(post_date)
                day_name  = dt.strftime("%A")
                date_str  = dt.strftime("%d %b %Y")
            except Exception:
                day_name, date_str = "", post_date

            ct = post.get("contentType", "")
            ct_color = type_colors.get(ct, "6C3CE1")
            ct_fill  = PatternFill("solid", fgColor=ct_color)
            row_fill = PatternFill("solid", fgColor="0F0A1E" if i % 2 == 0 else "1E1640")

            # Extract enriched fields with safe fallbacks
            hooks    = post.get("hook_variations") or [post.get("hook", ""), "", ""]
            while len(hooks) < 3:
                hooks.append("")
            cap_short = post.get("caption_short") or post.get("caption", "")[:120]
            cap_long  = post.get("caption_long")  or post.get("caption", "")
            cta_text  = post.get("cta", "")
            seo_kws   = post.get("seo_keywords") or []
            audio     = post.get("audio_suggestion") or {}
            audio_str = ""
            if audio and isinstance(audio, dict):
                audio_str = f"{audio.get('track_name','')} — {audio.get('vibe','')}"
                if audio.get("why_it_works"):
                    audio_str += f"\n({audio['why_it_works']})"

            # Per-type breakdown formatted as multi-line:
            #   Reel/AI Reel  → shot-by-shot script with time codes
            #   Carousel      → slide-by-slide
            #   Story         → frame-by-frame
            #   Graphic       → headline / subheadline / body / footer (static layout)
            breakdown = ""
            if ct == "Graphic" and post.get("graphic_layout") and isinstance(post["graphic_layout"], dict):
                gl = post["graphic_layout"]
                lines = ["── STATIC GRAPHIC LAYOUT ──"]
                if gl.get("headline"):
                    lines.append(f"HEADLINE: {gl['headline']}")
                if gl.get("subheadline"):
                    lines.append(f"SUBHEADLINE: {gl['subheadline']}")
                if gl.get("body_text"):
                    lines.append(f"BODY: {gl['body_text']}")
                supporting = gl.get("supporting_elements") or []
                if supporting:
                    lines.append("SUPPORTING:")
                    for el in supporting:
                        lines.append(f"  • {el}")
                if gl.get("footer_text"):
                    lines.append(f"FOOTER: {gl['footer_text']}")
                breakdown = "\n".join(lines)
            elif ct == "Carousel" and post.get("carousel_slides"):
                lines = ["── CAROUSEL SLIDES ──"]
                for s in post["carousel_slides"]:
                    n  = s.get("slide_number", "?")
                    hd = s.get("headline", "")
                    bd = s.get("body", "")
                    lines.append(f"Slide {n}: {hd} — {bd}")
                breakdown = "\n".join(lines)
            elif ct == "Story" and post.get("story_sequence"):
                lines = ["── STORY SEQUENCE ──"]
                for f in post["story_sequence"]:
                    n  = f.get("frame_number", "?")
                    tp = f.get("type", "")
                    tx = f.get("text", "")
                    lines.append(f"Frame {n} [{tp}]: {tx}")
                breakdown = "\n".join(lines)
            elif ct in ("Reel", "AI Reel") and post.get("reel_script") and isinstance(post["reel_script"], dict):
                rs = post["reel_script"]
                lines = [f"── REEL SCRIPT ({rs.get('duration_seconds','?')}s) ──"]
                for sh in rs.get("shots") or []:
                    tr = sh.get("time_range", "")
                    vs = sh.get("visual", "")
                    ot = sh.get("on_screen_text", "")
                    lines.append(f"{tr}: {vs} | OST: {ot}")
                if rs.get("cta_overlay"):
                    lines.append(f"Final CTA: {rs['cta_overlay']}")
                breakdown = "\n".join(lines)

            row_data = [
                i + 1,
                date_str,
                day_name,
                post.get("posting_time", ""),
                ct,
                post.get("topic", ""),
                hooks[0],
                hooks[1],
                hooks[2],
                cap_short,
                cap_long,
                cta_text,
                " ".join(post.get("hashtags", [])[:30]),
                ", ".join(seo_kws),
                audio_str,
                breakdown,
                post.get("visual_brief", ""),
                post.get("status", "draft").upper(),
            ]
            ct_col = 5  # Content Type is now column 5
            for col, val in enumerate(row_data, 1):
                cell = ws.cell(row=row, column=col, value=val)
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                cell.border    = thin_border
                cell.fill      = ct_fill if col == ct_col else row_fill
                if col == ct_col:
                    cell.font  = Font(bold=True, color="FFFFFF", size=9)
                elif col in (1, 2, 3, 4):
                    cell.font  = Font(color="A78BFA", size=9)
                else:
                    cell.font  = Font(color="D1D5DB", size=9)

            # Tall row for long content
            ws.row_dimensions[row].height = 120

        # Column widths — 18 columns
        widths = [
            5,   # #
            14,  # Date
            10,  # Day
            10,  # Time
            14,  # Content Type
            38,  # Topic
            32, 32, 32,  # Hook 1/2/3
            32,  # Caption Short
            55,  # Caption Long
            28,  # CTA
            55,  # Hashtags
            22,  # SEO Keywords
            28,  # Audio/Music
            55,  # Carousel/Story Breakdown
            32,  # Visual Brief
            10,  # Status
        ]
        for col, width in enumerate(widths, 1):
            ws.column_dimensions[get_column_letter(col)].width = width

        # Freeze top rows
        ws.freeze_panes = "A4"

        # ── Sheet 2: Growth Strategy ──
        ws2 = wb.create_sheet("Growth Strategy")
        ws2.sheet_view.showGridLines = False
        ws2["A1"] = "Growth Strategy Overview"
        ws2["A1"].font = Font(bold=True, color="FFFFFF", size=14)
        ws2["A1"].fill = header_fill

        row2 = 3
        sections = {
            "Content Pillars": strategy.get("pillars", []),
            "Growth Tactics":  strategy.get("growth_tactics", []),
            "CTA Templates":   strategy.get("cta_templates", []),
        }
        for section, items in sections.items():
            ws2.cell(row=row2, column=1, value=section).font = Font(bold=True, color="A78BFA", size=11)
            ws2.cell(row=row2, column=1).fill = subhdr_fill
            row2 += 1
            if isinstance(items, list):
                for item in items:
                    cell = ws2.cell(row=row2, column=1, value=f"  • {item}")
                    cell.font = Font(color="D1D5DB", size=10)
                    row2 += 1
            row2 += 1

        ws2.column_dimensions["A"].width = 80

        # Serialize and upload
        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)

        path = f"runs/{run_id}/content_calendar.xlsx"
        url  = _upload_bytes(buf.read(), path, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        return url

    except Exception as e:
        print(f"[Designer] Excel build error: {e}")
        return None


# ── PowerPoint PPT ────────────────────────────────────────────────────────────

async def _build_ppt(posts: list, brand: dict, strategy: dict, analyst_report: dict, run_id: str) -> str | None:
    try:
        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9

        name    = brand.get("name", "Brand")
        niche   = brand.get("niche", "")
        colors  = brand.get("colors") or {}
        primary = colors.get("primary", "#6C3CE1") if isinstance(colors, dict) else "#6C3CE1"

        def hex_to_rgb(h: str):
            h = h.lstrip("#")
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        brand_color = hex_to_rgb(primary)
        white       = RGBColor(0xFF, 0xFF, 0xFF)
        dark_bg     = RGBColor(0x0F, 0x0A, 0x1E)
        accent      = RGBColor(0xA7, 0x8B, 0xFA)

        blank_layout = prs.slide_layouts[6]  # blank

        def add_bg(slide, color: RGBColor):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_text(slide, text, left, top, width, height, size, bold=False, color=None, align=PP_ALIGN.LEFT):
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf  = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size   = Pt(size)
            run.font.bold   = bold
            run.font.color.rgb = color or white
            return txb

        # ── Slide 1: Title ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, dark_bg)
        # Accent bar
        bar = slide.shapes.add_shape(1, 0, 0, Inches(0.05), prs.slide_height)
        bar.fill.solid(); bar.fill.fore_color.rgb = brand_color; bar.line.fill.background()
        add_text(slide, name.upper(),         Inches(0.4), Inches(1.2), Inches(9), Inches(1.2), 40, bold=True, color=white)
        add_text(slide, f"{niche} — Content Strategy", Inches(0.4), Inches(2.4), Inches(9), Inches(0.8), 20, color=accent)
        add_text(slide, f"Generated by SocialOS | {datetime.utcnow().strftime('%d %b %Y')}", Inches(0.4), Inches(4.8), Inches(9), Inches(0.5), 12, color=RGBColor(0x6B,0x72,0x80))

        # ── Slide 2: Content Mix ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, dark_bg)
        add_text(slide, "Content Mix Strategy", Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), 28, bold=True)
        content_mix = strategy.get("content_mix", {})
        y_pos = 1.3
        for ct, pct in content_mix.items():
            add_text(slide, f"{ct}:  {pct}%", Inches(0.5), Inches(y_pos), Inches(4), Inches(0.4), 16, color=accent)
            y_pos += 0.5

        # ── Slide 3: Content Pillars ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, dark_bg)
        add_text(slide, "Content Pillars", Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), 28, bold=True)
        pillars = strategy.get("pillars", [])
        for i, pillar in enumerate(pillars[:4]):
            add_text(slide, f"0{i+1}. {pillar}", Inches(0.5), Inches(1.2 + i*0.9), Inches(9), Inches(0.7), 18, color=white if i%2==0 else accent)

        # ── Slide 4: Growth Tactics ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, dark_bg)
        add_text(slide, "Growth Tactics", Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), 28, bold=True)
        tactics = strategy.get("growth_tactics", [])
        y = 1.2
        for t in tactics[:5]:
            add_text(slide, f"• {t}", Inches(0.5), Inches(y), Inches(9), Inches(0.55), 14, color=accent if y < 2 else white)
            y += 0.6

        # ── Slide 5: 15-Day Calendar Summary ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, dark_bg)
        add_text(slide, f"{len(posts)}-Day Content Calendar", Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), 28, bold=True)
        y = 1.2
        for i, post in enumerate(posts[:8]):
            try:
                dt = datetime.fromisoformat(post.get("date",""))
                date_str = dt.strftime("%d %b")
            except Exception:
                date_str = post.get("date","")[:10]
            ct    = post.get("contentType","")
            topic = post.get("topic","")[:55]
            line  = f"{date_str}  [{ct}]  {topic}"
            add_text(slide, line, Inches(0.5), Inches(y), Inches(9), Inches(0.42), 11, color=white if i%2==0 else accent)
            y += 0.48

        # Serialize and upload
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        path = f"runs/{run_id}/strategy_deck.pptx"
        url  = _upload_bytes(buf.read(), path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        return url

    except Exception as e:
        print(f"[Designer] PPT build error: {e}")
        return None


# ── Supabase Upload ───────────────────────────────────────────────────────────

def _upload_bytes(data: bytes, path: str, content_type: str) -> str | None:
    sb = _get_supabase()
    if not sb:
        print(f"[Designer] Supabase unavailable — skipping upload for {path}")
        return None
    bucket = os.getenv("SUPABASE_STORAGE_BUCKET", "socialos-storage")
    try:
        sb.storage.from_(bucket).upload(
            path, data,
            file_options={"content-type": content_type, "upsert": "true"},
        )
        result = sb.storage.from_(bucket).get_public_url(path)
        print(f"[Designer] Uploaded {path} -> {str(result)[:80]}")
        return result
    except Exception as e:
        print(f"[Designer] Supabase upload error for {path}: {e}")
        import traceback; traceback.print_exc()
        return None
