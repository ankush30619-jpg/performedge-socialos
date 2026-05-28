"""
Performance Reporter Agent Node — 15-DAY COMPARISON REPORT
-----------------------------------------------------------
Generates a "Performance Report" PPT that compares the previous growth plan
(planned KPIs, content pillars, growth tactics) against actual results from
the last 15 days (live Meta engagement per published post).

Inputs from state:
  brand                  — full brand dict including decrypted igAccessToken
  analyst_report         — fresh Meta metrics from the analyst that just ran
  previous_strategy      — strategyJson from the previous completed run
  previous_analyst_report— analystReport from the previous run (baseline)
  previous_posts         — published posts in last 15d with igMediaId

Output:
  ppt_url                — Supabase URL of the comparison PPT
  _message               — one-line summary

Slides (6):
  1. Cover — "Performance Report — {brand}, Day {N}"
  2. Planned vs Actual KPIs (4-tile grid with ✓/✗ status)
  3. Top 5 posts that worked + why
  4. Bottom 5 posts that flopped + why
  5. Insights by content pillar (bar chart)
  6. Refined plan for next 15 days
"""
import asyncio
import io
import json
import os
from datetime import datetime
from typing import Any

import aiohttp
from openai import AsyncOpenAI

# Reuse the upload helper + Meta client signature from siblings
GRAPH_BASE = "https://graph.facebook.com/v21.0"
_oai = None
_supabase = None


def _get_oai():
    global _oai
    if _oai is not None:
        return _oai
    key = os.getenv("OPENAI_API_KEY", "")
    if key:
        _oai = AsyncOpenAI(api_key=key)
    return _oai


def _get_supabase():
    """Lazy singleton — same pattern as growth_planner._get_supabase."""
    global _supabase
    if _supabase is not None:
        return _supabase
    try:
        url = os.getenv("NEXT_PUBLIC_SUPABASE_URL", "")
        key = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
        if not url or not key:
            return None
        from supabase import create_client
        _supabase = create_client(url, key)
        return _supabase
    except Exception as e:
        print(f"[PerfReporter] Supabase init error: {e}")
        return None


def _upload_bytes(data: bytes, path: str, content_type: str) -> str | None:
    sb = _get_supabase()
    if not sb:
        print(f"[PerfReporter] Supabase unavailable — skipping upload for {path}")
        return None
    bucket = os.getenv("SUPABASE_STORAGE_BUCKET", "socialos-storage")
    try:
        sb.storage.from_(bucket).upload(
            path, data,
            file_options={"content-type": content_type, "upsert": "true"},
        )
        url = sb.storage.from_(bucket).get_public_url(path)
        print(f"[PerfReporter] Uploaded {path} -> {str(url)[:80]}")
        return url
    except Exception as e:
        print(f"[PerfReporter] Upload error for {path}: {e}")
        return None


# ── Meta per-post insights ────────────────────────────────────────────────────
async def _fetch_post_engagement(ig_media_id: str, access_token: str) -> dict:
    """Fetch live engagement for a single published post."""
    fields = "id,caption,media_type,timestamp,like_count,comments_count,permalink"
    url = f"{GRAPH_BASE}/{ig_media_id}?fields={fields}&access_token={access_token}"
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=10)) as r:
                data = await r.json()
        if "error" in data:
            return {"error": data["error"].get("message", "Meta error"), "id": ig_media_id}
        # Try to fetch the /insights endpoint for reach/saved/shares (optional)
        ins_fields = "reach,saved,shares,total_interactions"
        ins_url = f"{GRAPH_BASE}/{ig_media_id}/insights?metric={ins_fields}&access_token={access_token}"
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(ins_url, timeout=aiohttp.ClientTimeout(total=10)) as r:
                    ins = await r.json()
            for item in ins.get("data") or []:
                name = item.get("name")
                value = (item.get("values") or [{}])[0].get("value")
                if name and value is not None:
                    data[name] = value
        except Exception:
            pass  # insights endpoint can be flaky / require Business permissions
        return data
    except Exception as e:
        return {"error": str(e), "id": ig_media_id}


async def _enrich_posts_with_engagement(posts: list, access_token: str | None) -> list:
    """For each previous post, attach live Meta engagement numbers."""
    if not access_token or not posts:
        # Without a token we can still return the post shells with zero metrics
        return [{**p, "_metrics": {}} for p in posts]
    enriched = []
    # Run requests concurrently (cap at 8 parallel to stay polite)
    sem = asyncio.Semaphore(8)
    async def fetch_one(p):
        async with sem:
            m = await _fetch_post_engagement(p.get("igMediaId") or "", access_token)
        return {**p, "_metrics": m}
    results = await asyncio.gather(*(fetch_one(p) for p in posts), return_exceptions=False)
    enriched.extend(results)
    return enriched


def _score_post(p: dict) -> int:
    """Composite engagement score: likes + 3*comments + 5*saves."""
    m = p.get("_metrics") or {}
    likes    = int(m.get("like_count") or 0)
    comments = int(m.get("comments_count") or 0)
    saves    = int(m.get("saved") or 0)
    return likes + 3 * comments + 5 * saves


async def _gpt_post_insight(post: dict, kind: str) -> str:
    """Ask GPT for a 1-line 'why it worked / why it flopped' insight."""
    oai = _get_oai()
    if not oai:
        return ""
    m = post.get("_metrics") or {}
    summary = (
        f"Topic: {post.get('topic','')}\n"
        f"Content Type: {post.get('contentType','')}\n"
        f"Likes: {m.get('like_count', 0)} | Comments: {m.get('comments_count', 0)} | "
        f"Saves: {m.get('saved', 0)} | Reach: {m.get('reach', 'n/a')}"
    )
    try:
        resp = await oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a social media performance analyst. Return ONE short sentence (max 18 words) — no preamble."},
                {"role": "user", "content": f"Why did this post {kind}? Be specific and actionable.\n\n{summary}"},
            ],
            temperature=0.4,
            max_tokens=60,
        )
        return resp.choices[0].message.content.strip().strip('"').strip("'")
    except Exception:
        return ""


async def _gpt_refined_plan(brand: dict, previous_strategy: dict, current_analyst: dict, top_posts: list, bottom_posts: list) -> dict:
    """GPT generates: 3 keep-doing, 3 stop-doing, 3 new things to try."""
    oai = _get_oai()
    if not oai:
        return {"keep": [], "stop": [], "try": []}
    prev_pillars = (previous_strategy or {}).get("pillars") or []
    prev_mix     = (previous_strategy or {}).get("content_mix") or {}
    name         = brand.get("name", "the brand")
    niche        = brand.get("niche", "")
    top_topics   = [p.get("topic", "") for p in top_posts[:3]]
    flop_topics  = [p.get("topic", "") for p in bottom_posts[:3]]

    user_prompt = (
        f"Brand: {name} ({niche})\n"
        f"Previous content pillars: {prev_pillars}\n"
        f"Previous content mix: {prev_mix}\n"
        f"Top performing topics last 15d: {top_topics}\n"
        f"Bottom performing topics: {flop_topics}\n"
        f"Current Meta engagement rate: {current_analyst.get('avgEngagementRate', 'n/a')}%\n\n"
        f"Return JSON: {{ \"keep\": [3 strings], \"stop\": [3 strings], \"try\": [3 strings] }}\n"
        f"Each string is a concrete action (8-14 words). Reference specifics, not generic advice."
    )
    try:
        resp = await oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a brand growth strategist. Return JSON only — no prose."},
                {"role": "user", "content": user_prompt},
            ],
            response_format={"type": "json_object"},
            temperature=0.7,
            max_tokens=600,
        )
        data = json.loads(resp.choices[0].message.content)
        # Normalise — ensure each list is exactly 3
        for k in ("keep", "stop", "try"):
            lst = data.get(k) or []
            if not isinstance(lst, list):
                lst = [str(lst)]
            data[k] = (lst + ["—"] * 3)[:3]
        return data
    except Exception as e:
        print(f"[PerfReporter] GPT refined plan error: {e}")
        return {"keep": [], "stop": [], "try": []}


# ── Main node ─────────────────────────────────────────────────────────────────
async def performance_reporter_node(state: dict, event_queue: asyncio.Queue) -> dict:
    brand                  = state.get("brand") or {}
    analyst_report         = state.get("analyst_report") or {}
    previous_strategy      = state.get("previous_strategy") or {}
    previous_analyst_report = state.get("previous_analyst_report") or {}
    previous_posts         = state.get("previous_posts") or []
    run_id                 = state.get("run_id") or "perf"

    name           = brand.get("name", "Brand")
    ig_token       = brand.get("igAccessToken")
    posts_count    = len(previous_posts)

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "performanceReporter",
        "message": f"Building performance report — {posts_count} published posts to analyse…",
    })

    # 1. Enrich the previous posts with live Meta engagement
    enriched_posts = await _enrich_posts_with_engagement(previous_posts, ig_token)

    # 2. Sort by composite engagement score
    scored = sorted(enriched_posts, key=_score_post, reverse=True)
    top_posts    = scored[:5]
    bottom_posts = list(reversed(scored))[:5] if len(scored) > 5 else scored[5:]

    # 3. Ask GPT for "why" insights on top/bottom posts (concurrent)
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "performanceReporter",
        "message": "Asking GPT for per-post insights…",
    })
    top_insights, bottom_insights = await asyncio.gather(
        asyncio.gather(*(_gpt_post_insight(p, "perform so well") for p in top_posts)),
        asyncio.gather(*(_gpt_post_insight(p, "underperform")    for p in bottom_posts)),
    )
    for p, ins in zip(top_posts,    top_insights):    p["_insight"] = ins
    for p, ins in zip(bottom_posts, bottom_insights): p["_insight"] = ins

    # 4. Refined plan
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "performanceReporter",
        "message": "Generating refined plan for next 15 days…",
    })
    refined_plan = await _gpt_refined_plan(brand, previous_strategy, analyst_report, top_posts, bottom_posts)

    # 5. Build the PPT
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "performanceReporter",
        "message": "Composing PPT…",
    })
    ppt_url = await _build_performance_ppt(
        brand=brand,
        analyst_report=analyst_report,
        previous_strategy=previous_strategy,
        previous_analyst_report=previous_analyst_report,
        top_posts=top_posts,
        bottom_posts=bottom_posts,
        refined_plan=refined_plan,
        days=15,
        run_id=run_id,
    )

    return {
        "ppt_url": ppt_url,
        "posts_generated": 0,  # this mode doesn't create new posts
        "agent_statuses": {
            **(state.get("agent_statuses") or {}),
            "performanceReporter": {"status": "completed", "message": f"Performance report ready — {posts_count} posts analysed"},
        },
        "_message": f"Performance report ready — {posts_count} posts analysed, {len(top_posts)} winners highlighted",
    }


# ── PPT Builder ───────────────────────────────────────────────────────────────
async def _build_performance_ppt(
    brand: dict,
    analyst_report: dict,
    previous_strategy: dict,
    previous_analyst_report: dict,
    top_posts: list,
    bottom_posts: list,
    refined_plan: dict,
    days: int,
    run_id: str,
) -> str | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.625)

        name    = brand.get("name", "Brand")
        niche   = brand.get("niche", "")
        colors  = brand.get("colors") or {}
        primary = colors.get("primary", "#6C3CE1") if isinstance(colors, dict) else "#6C3CE1"

        def hex_to_rgb(h):
            h = h.lstrip("#")
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        brand_color = hex_to_rgb(primary)
        white  = RGBColor(0xFF, 0xFF, 0xFF)
        dark   = RGBColor(0x0F, 0x0A, 0x1E)
        card_dark = RGBColor(0x1E, 0x16, 0x40)
        muted  = RGBColor(0xA7, 0x8B, 0xFA)
        green  = RGBColor(0x10, 0xB9, 0x81)
        red    = RGBColor(0xEF, 0x44, 0x44)
        yellow = RGBColor(0xF5, 0x9E, 0x0B)

        blank = prs.slide_layouts[6]

        def bg(slide, color):
            f = slide.background.fill; f.solid(); f.fore_color.rgb = color

        def txt(slide, text, l, t, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
            tb = slide.shapes.add_textbox(l, t, w, h)
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = align
            r  = p.add_run(); r.text = str(text)
            r.font.size  = Pt(size); r.font.bold = bold
            r.font.color.rgb = color or white
            r.font.italic = italic

        def bar(slide, l, t, w, h, color):
            s = slide.shapes.add_shape(1, l, t, w, h)
            s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

        # ── Slide 1: Cover ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(0.06), prs.slide_height, brand_color)
        txt(s, "PERFORMANCE REPORT", Inches(0.4), Inches(0.6), Inches(9), Inches(0.6), 14, color=muted)
        txt(s, name.upper(),         Inches(0.4), Inches(1.2), Inches(9), Inches(1.2), 40, bold=True)
        txt(s, f"{niche} — Day {days} Review",      Inches(0.4), Inches(2.5), Inches(9), Inches(0.6), 18, color=muted)
        txt(s, f"Generated by SocialOS · {datetime.utcnow().strftime('%d %b %Y')}",
            Inches(0.4), Inches(4.8), Inches(9), Inches(0.4), 11, color=RGBColor(0x6B,0x72,0x80))

        # ── Slide 2: Planned vs Actual KPIs (4-tile grid) ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(10), Inches(0.08), brand_color)
        txt(s, "Planned vs Actual — Last 15 Days", Inches(0.4), Inches(0.3), Inches(9), Inches(0.5), 22, bold=True)
        txt(s, "How we did against the plan", Inches(0.4), Inches(0.85), Inches(9), Inches(0.35), 12, color=muted)

        # Pull planned KPI targets from previous strategy
        prev_kpi    = (previous_strategy or {}).get("kpi_targets_90day") or {}
        prev_mix    = (previous_strategy or {}).get("content_mix") or {}
        target_followers = (prev_kpi.get("followers_target")
                            or (previous_analyst_report or {}).get("followerCount", 0))
        target_er        = float(prev_kpi.get("engagement_rate_target", 2.5) or 2.5)
        target_reach     = int(prev_kpi.get("reach_per_post_target", 100) or 100)
        planned_posts    = sum(int(v) for v in prev_mix.values()) if prev_mix else 15

        actual_followers = int(analyst_report.get("followerCount", 0) or 0)
        actual_er        = float(analyst_report.get("avgEngagementRate", 0) or 0)
        actual_reach     = int(analyst_report.get("avgReach", 0) or 0)
        actual_posts     = len(top_posts) + len(bottom_posts)  # all enriched

        def tile_color(actual, target, lower_is_worse=True):
            if not target:
                return yellow
            ratio = actual / target if lower_is_worse else target / max(actual, 1)
            if ratio >= 1.0:  return green
            if ratio >= 0.8:  return yellow
            return red

        def delta_str(actual, target):
            if not target:
                return ""
            diff = actual - target
            sign = "+" if diff >= 0 else ""
            return f"{sign}{diff} ({sign}{(diff/target*100):.0f}%)"

        tiles = [
            ("Followers",       f"{actual_followers:,}", f"Target: {int(target_followers):,}", delta_str(actual_followers, target_followers), tile_color(actual_followers, target_followers)),
            ("Engagement Rate", f"{actual_er:.2f}%",     f"Target: {target_er:.2f}%",          delta_str(actual_er, target_er),                tile_color(actual_er, target_er)),
            ("Reach / Post",    f"{actual_reach:,}",     f"Target: {target_reach:,}",          delta_str(actual_reach, target_reach),          tile_color(actual_reach, target_reach)),
            ("Posts Published", f"{actual_posts}",       f"Planned: {planned_posts}",          delta_str(actual_posts, planned_posts),         tile_color(actual_posts, planned_posts)),
        ]
        tile_w = Inches(2.2); tile_h = Inches(2.4); gap = Inches(0.15); start_x = Inches(0.4); y = Inches(1.5)
        for i, (label, value, target_txt, delta, color) in enumerate(tiles):
            x = start_x + (tile_w + gap) * i
            bar(s, x, y, tile_w, tile_h, card_dark)
            bar(s, x, y, tile_w, Inches(0.1), color)
            check = "✅" if color == green else ("⚠️" if color == yellow else "❌")
            txt(s, label,      x + Inches(0.15), y + Inches(0.2), tile_w - Inches(0.3), Inches(0.4), 11, color=muted)
            txt(s, value,      x + Inches(0.15), y + Inches(0.65), tile_w - Inches(0.3), Inches(0.8), 28, bold=True, color=color)
            txt(s, target_txt, x + Inches(0.15), y + Inches(1.55), tile_w - Inches(0.3), Inches(0.35), 10, color=white)
            txt(s, f"{check} {delta}", x + Inches(0.15), y + Inches(1.9), tile_w - Inches(0.3), Inches(0.35), 11, bold=True, color=color)

        # ── Slide 3: Top 5 Posts That Worked ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(10), Inches(0.08), green)
        txt(s, "✅ Top 5 Posts That Worked", Inches(0.4), Inches(0.3), Inches(9), Inches(0.5), 22, bold=True, color=green)
        txt(s, "Highest engagement — replicate this style next 15 days", Inches(0.4), Inches(0.85), Inches(9), Inches(0.35), 12, color=muted)

        if not top_posts:
            txt(s, "Publish posts via the Calendar to populate this report.", Inches(0.4), Inches(2.5), Inches(9), Inches(0.5), 16, italic=True, color=muted, align=PP_ALIGN.CENTER)
        else:
            row_y = Inches(1.4)
            for i, p in enumerate(top_posts[:5]):
                m = p.get("_metrics") or {}
                row_h = Inches(0.65)
                bar(s, Inches(0.4), row_y, Inches(9.2), row_h, card_dark)
                bar(s, Inches(0.4), row_y, Inches(0.06), row_h, green)
                # Topic + content type
                ct = p.get("contentType", "")
                txt(s, f"#{i+1}  {p.get('topic','')[:55]}", Inches(0.6), row_y + Inches(0.05), Inches(5.8), Inches(0.32), 11, bold=True)
                txt(s, f"[{ct}]  {p.get('_insight','')[:80]}", Inches(0.6), row_y + Inches(0.36), Inches(5.8), Inches(0.28), 9, color=muted, italic=True)
                # Metrics
                metric_txt = f"♥ {m.get('like_count', 0)}   💬 {m.get('comments_count', 0)}   🔖 {m.get('saved', 0)}"
                if m.get('reach'):
                    metric_txt += f"   👁 {m['reach']}"
                txt(s, metric_txt, Inches(6.5), row_y + Inches(0.18), Inches(3.0), Inches(0.4), 11, bold=True, color=green, align=PP_ALIGN.RIGHT)
                row_y += Inches(0.75)

        # ── Slide 4: Bottom 5 Posts That Flopped ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(10), Inches(0.08), red)
        txt(s, "❌ Bottom 5 Posts That Flopped", Inches(0.4), Inches(0.3), Inches(9), Inches(0.5), 22, bold=True, color=red)
        txt(s, "Low engagement — stop or rework these formats", Inches(0.4), Inches(0.85), Inches(9), Inches(0.35), 12, color=muted)

        if not bottom_posts:
            txt(s, "Not enough data yet — needs at least 6 published posts.", Inches(0.4), Inches(2.5), Inches(9), Inches(0.5), 16, italic=True, color=muted, align=PP_ALIGN.CENTER)
        else:
            row_y = Inches(1.4)
            for i, p in enumerate(bottom_posts[:5]):
                m = p.get("_metrics") or {}
                row_h = Inches(0.65)
                bar(s, Inches(0.4), row_y, Inches(9.2), row_h, card_dark)
                bar(s, Inches(0.4), row_y, Inches(0.06), row_h, red)
                ct = p.get("contentType", "")
                txt(s, f"#{i+1}  {p.get('topic','')[:55]}", Inches(0.6), row_y + Inches(0.05), Inches(5.8), Inches(0.32), 11, bold=True)
                txt(s, f"[{ct}]  {p.get('_insight','')[:80]}", Inches(0.6), row_y + Inches(0.36), Inches(5.8), Inches(0.28), 9, color=muted, italic=True)
                metric_txt = f"♥ {m.get('like_count', 0)}   💬 {m.get('comments_count', 0)}   🔖 {m.get('saved', 0)}"
                txt(s, metric_txt, Inches(6.5), row_y + Inches(0.18), Inches(3.0), Inches(0.4), 11, bold=True, color=red, align=PP_ALIGN.RIGHT)
                row_y += Inches(0.75)

        # ── Slide 5: Insights by Content Pillar ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(10), Inches(0.08), brand_color)
        txt(s, "Performance by Content Pillar", Inches(0.4), Inches(0.3), Inches(9), Inches(0.5), 22, bold=True)
        txt(s, "Average engagement by content type — double down on the winners", Inches(0.4), Inches(0.85), Inches(9), Inches(0.35), 12, color=muted)

        # Group by content type
        type_engagement: dict[str, list[int]] = {}
        all_enriched = top_posts + bottom_posts
        for p in all_enriched:
            ct = p.get("contentType", "Other")
            score = _score_post(p)
            type_engagement.setdefault(ct, []).append(score)
        # Compute averages
        type_avg = sorted(
            [(ct, sum(scores)/len(scores)) for ct, scores in type_engagement.items()],
            key=lambda x: x[1], reverse=True
        )

        if not type_avg:
            txt(s, "Need at least 1 published post per content type to draw chart.", Inches(0.4), Inches(2.5), Inches(9), Inches(0.5), 14, italic=True, color=muted, align=PP_ALIGN.CENTER)
        else:
            max_score = max(score for _, score in type_avg) or 1
            row_y = Inches(1.5)
            type_colors = {
                "Reel":     RGBColor(0x3B, 0x82, 0xF6),
                "AI Reel":  RGBColor(0xEC, 0x48, 0x99),
                "Carousel": RGBColor(0x10, 0xB9, 0x81),
                "Graphic":  RGBColor(0xA7, 0x8B, 0xFA),
                "Story":    RGBColor(0xF5, 0x9E, 0x0B),
            }
            for ct, avg in type_avg[:6]:
                color = type_colors.get(ct, brand_color)
                bar_w = Inches(0.4 + (avg / max_score) * 6.5)
                txt(s, ct, Inches(0.4), row_y, Inches(1.8), Inches(0.4), 12, bold=True, color=color)
                bar(s, Inches(2.4), row_y + Inches(0.05), bar_w, Inches(0.32), color)
                txt(s, f"{avg:.0f} pts", Inches(2.4) + bar_w + Inches(0.1), row_y, Inches(2), Inches(0.4), 11, color=white)
                row_y += Inches(0.55)

        # ── Slide 6: Refined Plan for Next 15 Days ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(10), Inches(0.08), brand_color)
        txt(s, "Refined Plan — Next 15 Days", Inches(0.4), Inches(0.3), Inches(9), Inches(0.5), 22, bold=True)
        txt(s, "Adjusted based on the wins and misses above", Inches(0.4), Inches(0.85), Inches(9), Inches(0.35), 12, color=muted)

        cols = [
            ("✅ KEEP DOING", refined_plan.get("keep", []), green,   Inches(0.4)),
            ("🛑 STOP DOING", refined_plan.get("stop", []), red,     Inches(3.6)),
            ("⚡ TRY NEW",     refined_plan.get("try",  []), yellow, Inches(6.8)),
        ]
        col_w = Inches(3.0); col_h = Inches(3.4); col_y = Inches(1.5)
        for header, items, color, x in cols:
            bar(s, x, col_y, col_w, col_h, card_dark)
            bar(s, x, col_y, col_w, Inches(0.08), color)
            txt(s, header, x + Inches(0.2), col_y + Inches(0.15), col_w - Inches(0.4), Inches(0.4), 13, bold=True, color=color)
            iy = col_y + Inches(0.7)
            for item in items[:3]:
                txt(s, f"• {item}", x + Inches(0.2), iy, col_w - Inches(0.4), Inches(0.9), 10, color=white)
                iy += Inches(0.9)
        txt(s, "Auto-applied to the next pipeline run", Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.4), 10, color=muted, italic=True, align=PP_ALIGN.CENTER)

        # ── Serialize + upload ──
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        path = f"runs/{run_id}/performance_report.pptx"
        url = _upload_bytes(buf.read(), path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        return url

    except Exception as e:
        print(f"[PerfReporter] PPT build error: {e}")
        import traceback; traceback.print_exc()
        return None
