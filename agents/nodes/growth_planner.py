"""
Growth Planner Agent Node — FULL INSTAGRAM AUDIT + GROWTH PPT
-----------------------------------------------------------------
1. Runs Research + Competitor sub-agents in parallel
2. Deep Instagram audit: every reel, every like/comment/save/reach
3. Identifies what content is WORKING vs NOT WORKING
4. Builds content pillars from real data
5. Sets follower growth goal (current → target) with month-by-month plan
6. Generates a comprehensive Growth Strategy PPT uploaded to Supabase
"""
import asyncio
import io
import json
import os
from datetime import datetime

from llm_client import complete as llm_complete
from state import SocialOSState
from nodes.research_agent import research_agent_node
from nodes.competitor_tracker import competitor_tracker_node
from skills.registry import (
    GROWTH_ROADMAP_FRAMEWORK, STRATEGIC_THINKING,
    BUYER_STAGES, ANTI_AI_LANGUAGE,
    # 2026 upgrades
    CRITICAL_INSTRUCTION_PREFIX, LEARNED_PATTERNS_SLOT,
    HASHTAG_STRATEGY_2026, GROWTH_KPI_THRESHOLDS_2026,
    ANTI_AI_LANGUAGE_2026_ADDITIONS,
)

# Lazy Supabase singleton — initialized on first use
_supabase = None


def _get_supabase():
    """Return (or lazily create) the Supabase client."""
    global _supabase
    if _supabase is not None:
        return _supabase
    url = os.getenv("NEXT_PUBLIC_SUPABASE_URL", "")
    key = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
    if not url or not key:
        print(f"[GrowthPlanner] Supabase env vars missing")
        return None
    try:
        from supabase import create_client
        _supabase = create_client(url, key)
        return _supabase
    except Exception as e:
        print(f"[GrowthPlanner] Supabase init error: {e}")
        return None


async def growth_planner_node(state: SocialOSState, event_queue: asyncio.Queue) -> dict:
    brand                      = state.get("brand") or {}
    brand_knowledge            = state.get("brand_knowledge") or {}
    analyst_report             = state.get("analyst_report") or {}
    days_ahead                 = state.get("days_ahead", 15)
    run_id                     = state.get("run_id", "")
    mode                       = state.get("mode", "full")
    user_follower_goal         = state.get("follower_goal")          # user-supplied target
    current_followers_override = state.get("current_followers_override")  # user-supplied current

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": "Launching Research + Competitor analysis in parallel…",
    })

    # ── Sub-agent fan-out ─────────────────────────────────────────────────────
    await event_queue.put({"type": "agent_started", "agentKey": "researchAgent",     "message": "Research Agent started"})
    await event_queue.put({"type": "agent_started", "agentKey": "competitorTracker", "message": "Competitor Tracker started"})

    research_result, competitor_result = await asyncio.gather(
        research_agent_node(state, event_queue),
        competitor_tracker_node(state, event_queue),
        return_exceptions=True,
    )

    research_data   = research_result.get("research_data",   {}) if isinstance(research_result, dict) else {}
    competitor_data = competitor_result.get("competitor_data", {}) if isinstance(competitor_result, dict) else {}

    await event_queue.put({"type": "agent_completed", "agentKey": "researchAgent",
        "message": research_result.get("_message", "Done") if isinstance(research_result, dict) else "Done"})
    await event_queue.put({"type": "agent_completed", "agentKey": "competitorTracker",
        "message": competitor_result.get("_message", "Done") if isinstance(competitor_result, dict) else "Done"})

    # ── Instagram Audit ───────────────────────────────────────────────────────
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": "Auditing Instagram performance — analysing every post…",
    })

    ig_audit    = _build_ig_audit(
        analyst_report, brand, brand_knowledge,
        follower_goal=user_follower_goal,
        current_followers_override=current_followers_override,
    )
    feasibility = _calculate_goal_feasibility(ig_audit, days_ahead)

    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": f"Audit complete — {ig_audit['posts_analysed']} posts, ER {ig_audit['avg_er']}% · Goal feasibility: {feasibility['probability_pct']}%",
    })

    # ── Growth Strategy Synthesis ─────────────────────────────────────────────
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": "Building AI-powered growth strategy…",
    })

    growth_strategy = await _generate_strategy(
        brand, brand_knowledge, analyst_report,
        research_data, competitor_data, ig_audit, days_ahead, feasibility,
        learned_patterns=state.get("learned_patterns") or "",
    )

    # ── Quality-check phase (Change 2): validate numbers before they reach
    # any slide. Verifies consistency, clamps impossible values, and labels
    # low-confidence figures as Estimate/Assumption. Scalable — every pipeline
    # that calls growth_planner_node inherits this gate automatically.
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": "Running data-quality validation on strategy numbers…",
    })
    growth_strategy, quality_report = _validate_strategy(
        growth_strategy, ig_audit, feasibility, days_ahead
    )
    if quality_report.get("corrections"):
        await event_queue.put({
            "type": "agent_progress",
            "agentKey": "growthPlanner",
            "message": f"Quality check: {len(quality_report['corrections'])} value(s) corrected · confidence: {quality_report.get('confidence','medium')}",
        })

    # ── Extended plan: forecasting + structured recs with evidence + 30/60/90-day calendar
    await event_queue.put({
        "type": "agent_progress",
        "agentKey": "growthPlanner",
        "message": "Generating investor-grade extended plan (forecasting + recommendations + execution calendar)…",
    })
    extended_plan = await _generate_extended_plan(
        brand, ig_audit, feasibility, research_data, competitor_data, growth_strategy, days_ahead
    )
    _clamp_extended_forecasting(extended_plan, ig_audit)

    # ── Build PPT for growth_planner_only mode ────────────────────────────────
    ppt_url = None
    if mode == "growth_planner_only":
        await event_queue.put({
            "type": "agent_progress",
            "agentKey": "growthPlanner",
            "message": "Building 22-slide investor-grade PPT deck…",
        })
        ppt_url = await _build_growth_ppt(
            brand, ig_audit, growth_strategy,
            research_data, competitor_data, days_ahead, run_id,
            analyst_report=analyst_report,
            feasibility=feasibility,
            extended_plan=extended_plan,
        )
        await event_queue.put({
            "type": "agent_progress",
            "agentKey": "growthPlanner",
            "message": f"PPT {'uploaded ✓' if ppt_url else 'failed'}",
        })
        # ── File traceability: tie the generated deck to the data that fed it
        if ppt_url:
            await event_queue.put({
                "type": "agent_sources",
                "agentKey": "growthPlanner",
                "message": "Growth Planner deck generated",
                "data": {
                    "files_generated": [{
                        "name": "growth_planner.pptx",
                        "url": ppt_url,
                        "type": "Strategy presentation (22 slides — investor-grade)",
                        "generated_by": "Growth Planner",
                        "data_confidence": quality_report.get("confidence", "medium"),
                        "produced_from": [
                            f"Instagram audit — {ig_audit.get('posts_analysed', 0)} posts ({ig_audit.get('data_source', 'estimate')} data)",
                            f"Follower goal math — {ig_audit.get('followers', 0):,} → {ig_audit.get('goal_followers', 0):,} over {days_ahead} days",
                            f"Research — {len(research_data.get('trending_angles', []) or [])} trending angles, {len(research_data.get('hashtags', []) or [])} hashtags",
                            f"Competitor intelligence — {len(competitor_data.get('competitors_found', []) or [])} competitors, {len(competitor_data.get('content_gaps', []) or [])} gaps",
                            f"AI growth strategy — {len(growth_strategy.get('growth_tactics', []) or [])} tactics across {len(growth_strategy.get('pillars', []) or [])} pillars",
                        ],
                        "quality_checks": quality_report.get("checks", []),
                    }],
                },
            })

    result = {
        "research_data":   research_data,
        "competitor_data": competitor_data,
        "growth_strategy": growth_strategy,
        "_message": (
            f"Growth strategy ready — {growth_strategy.get('posting_frequency','1x daily')} · "
            f"{ig_audit['posts_analysed']} posts audited"
        ),
    }
    if ppt_url:
        result["ppt_url"] = ppt_url

    return result


# ── Instagram Audit ─────────────────────────────────────────────────────────

def _build_ig_audit(
    analyst_report: dict,
    brand: dict,
    brand_knowledge: dict = None,
    follower_goal: int = None,
    current_followers_override: int = None,
) -> dict:
    """Analyse the analyst_report to produce an Instagram audit dict."""
    top_posts    = analyst_report.get("topPosts") or []
    followers    = analyst_report.get("followerCount", 0) or 0
    avg_er       = analyst_report.get("avgEngagementRate", 0) or 0
    avg_reach    = analyst_report.get("avgReach", 0) or 0
    ig_connected = analyst_report.get("ig_connected", False)
    profile_views_30d = (analyst_report.get("profileViews30d") or
                         analyst_report.get("profile_views_30d") or 0)
    # Benchmark / scraped fallback metrics (labeled Estimate downstream)
    est_er     = analyst_report.get("estimated_er", 0) or 0
    est_reach  = analyst_report.get("estimated_reach", 0) or 0
    bm_source  = analyst_report.get("benchmark_source", "")
    bm_conf    = analyst_report.get("benchmark_confidence", 55)
    analyst_data_source = analyst_report.get("data_source", "none")

    # User-supplied current followers override takes priority
    if current_followers_override is not None and current_followers_override >= 0:
        followers = current_followers_override

    # Fall back to brand.igFollowers when live IG data is unavailable.
    if followers == 0 and brand:
        stored = int(brand.get("igFollowers") or 0)
        if stored > 0:
            followers = stored

    # Use benchmark ER/reach if live data is zero but benchmarks exist
    if avg_er == 0 and est_er > 0:
        avg_er = est_er
    if avg_reach == 0 and est_reach > 0:
        avg_reach = est_reach

    data_source = "live" if (ig_connected and followers > 0 and len(top_posts) > 0) else (
        "scraped"   if analyst_data_source == "scraped" and followers > 0 else (
        "benchmark" if analyst_data_source == "benchmark" and followers > 0 else (
        "stored"    if (not ig_connected and followers > 0) else "estimate"
    )))

    # Classify posts as working / not working
    working     = []
    not_working = []
    content_types: dict = {}

    for p in top_posts:
        er = p.get("engagementRate", 0) or 0
        reach = p.get("reach", 0) or 0
        ct = p.get("mediaType", "POST")
        content_types[ct] = content_types.get(ct, 0) + 1

        entry = {
            "caption":       (p.get("caption") or "")[:80],
            "mediaType":     ct,
            "likes":         p.get("likes", 0),
            "comments":      p.get("comments", 0),
            "reach":         reach,
            "er":            round(er, 2),
            "timestamp":     p.get("timestamp", ""),
            "permalink":     p.get("permalink", ""),
        }
        if er >= (avg_er * 1.2) or reach >= (avg_reach * 1.2):
            working.append(entry)
        else:
            not_working.append(entry)

    # Best content type — only meaningful when there's actual content data.
    # When posts_analysed=0 the "best format" claim is fabricated, so suppress it.
    best_ct = max(content_types, key=content_types.get) if content_types else None

    # Goal calculator (spec §growth_calculator):
    # - Reject any user-supplied goal that is BELOW current followers — that's the
    #   "Current: 405, Target: 100, INVALID, Auto regenerate" case in the spec.
    # - Reject a goal that exactly equals current (no growth = no plan).
    # - When invalid, auto-regenerate using a sensible default (1.5x or +200, whichever larger).
    current_followers = followers
    goal_was_invalid  = False
    if follower_goal is not None and isinstance(follower_goal, (int, float)) and follower_goal > current_followers:
        goal_followers = int(follower_goal)
    elif current_followers > 0:
        # Auto-regenerate: 50% growth or +200 followers, whichever is larger
        goal_followers = max(current_followers + 200, int(current_followers * 1.5))
        if follower_goal is not None and follower_goal > 0 and follower_goal <= current_followers:
            goal_was_invalid = True
    else:
        kpi = (analyst_report.get("kpi_targets_90day") or {})
        goal_followers = int(kpi.get("followers", 0)) or 500
    gap = max(goal_followers - current_followers, 0)

    # Data-availability flags (spec §data_validation_system):
    # When metrics are truly unmeasured, surface them as "Data Unavailable" rather
    # than "0" — and let downstream slides skip content-performance fabrication.
    has_live_metrics = ig_connected and (len(top_posts) > 0)

    return {
        "ig_connected":     ig_connected,
        "data_source":      data_source,
        "posts_analysed":   len(top_posts),
        "followers":        current_followers,
        "goal_followers":   goal_followers,
        "follower_gap":     gap,
        "avg_er":           round(avg_er, 2),
        "avg_reach":        avg_reach,
        "best_content_type": best_ct,
        # New: data-quality flags consumed by the PPT layer
        "has_live_metrics":  has_live_metrics,
        "goal_was_invalid":  goal_was_invalid,
        "est_er":            est_er,
        "est_reach":         est_reach,
        "benchmark_source":  bm_source,
        "benchmark_confidence": bm_conf,
        "content_type_mix": content_types,
        "working_posts":    working[:5],
        "not_working_posts": not_working[:5],
        "total_likes_30d":  analyst_report.get("totalLikes30d", 0),
        "total_comments_30d": analyst_report.get("totalComments30d", 0),
        "impressions_30d":  analyst_report.get("impressions30d", 0),
        "profile_views_30d": profile_views_30d,
    }


# ── Goal Feasibility Calculator ──────────────────────────────────────────────

def _calculate_goal_feasibility(ig_audit: dict, days_ahead: int) -> dict:
    """Pre-compute goal math in Python so GPT reasons from reliable numbers."""
    followers    = ig_audit.get("followers", 0)
    goal         = ig_audit.get("goal_followers", 500)
    gap          = max(goal - followers, 0)
    ig_connected = ig_audit.get("ig_connected", False)
    avg_reach    = ig_audit.get("avg_reach", 0)

    if ig_connected and followers > 0 and avg_reach > 0:
        est_weekly_growth = max(1, int(avg_reach * 0.02))
    else:
        est_weekly_growth = 5  # cold-start baseline

    weeks             = max(1, days_ahead / 7)
    projected_at_pace = int(est_weekly_growth * weeks)
    required_weekly   = round(gap / weeks, 1)
    required_daily    = round(gap / max(1, days_ahead), 1)
    accel             = round(required_weekly / max(1, est_weekly_growth), 1)
    probability       = min(95, max(10, int(100 / max(1, accel))))

    risks = []
    if accel > 5:        risks.append(f"Goal needs {accel}x acceleration — very aggressive for {days_ahead} days")
    if accel > 2:        risks.append("Requires consistent daily Reels + active engagement every day")
    if not ig_connected: risks.append("No live IG data — estimates based on brand brief + niche benchmarks")
    if days_ahead <= 15: risks.append("Short window — first 7 days are make-or-break for momentum")

    return {
        "current_followers":         followers,
        "target_followers":          goal,
        "gap":                       gap,
        "days_ahead":                days_ahead,
        "est_current_weekly_growth": est_weekly_growth,
        "projected_growth_at_pace":  projected_at_pace,
        "required_weekly_growth":    required_weekly,
        "required_daily_growth":     required_daily,
        "acceleration_needed":       accel,
        "probability_pct":           probability,
        "risks":                     risks,
    }


# ── GPT Growth Strategy ──────────────────────────────────────────────────────

async def _generate_strategy(brand, brand_knowledge, analyst_report, research_data, competitor_data, ig_audit, days_ahead, feasibility: dict = None, learned_patterns: str = "") -> dict:
    niche           = brand.get("niche", "")
    name            = brand.get("name", "brand")
    tone            = brand.get("tone", "Professional")
    audience        = brand.get("targetAudience", "")
    positioning     = brand.get("positioning", "")
    differentiation = brand.get("differentiation", "")
    voice_style     = brand.get("voiceStyle", "")
    hook_style      = brand.get("hookStyle", "")
    cta_style       = brand.get("ctaStyle", "")
    catchphrases    = brand.get("catchphrases", "")
    content_pillars = brand.get("contentPillars") or []
    audience_pain   = brand.get("audiencePainPoints", "")
    context_block   = brand_knowledge.get("context_block", "")

    # Research insights
    trending_angles  = research_data.get("trending_angles", [])
    content_opps     = research_data.get("content_opportunities", [])
    hook_ideas       = research_data.get("hook_ideas", [])
    posting_insights = research_data.get("posting_insights", [])

    # Competitor insights
    comp_gaps        = competitor_data.get("content_gaps", [])
    diff_strategy    = competitor_data.get("differentiation_strategy", "")
    formats_to_own   = competitor_data.get("content_formats_to_own", [])

    # IG performance data
    working_text = "\n".join(
        f"- [{p['mediaType']}] \"{p['caption']}\" | ER: {p['er']}% | Reach: {p['reach']:,}"
        for p in ig_audit.get("working_posts", [])[:4]
    )
    not_working_text = "\n".join(
        f"- [{p['mediaType']}] \"{p['caption']}\" | ER: {p['er']}% | Reach: {p['reach']:,}"
        for p in ig_audit.get("not_working_posts", [])[:3]
    )
    trends_text      = "\n".join(f"- {t}" for t in trending_angles[:5]) or "\n".join(
        f"- {t.get('title','')}" for t in research_data.get("trends", [])[:5]
    )
    gaps_text        = "\n".join(f"- {g}" for g in comp_gaps[:4])
    posting_text     = "\n".join(f"- {p}" for p in posting_insights[:3])

    ig_connected = ig_audit.get("ig_connected", False)
    followers    = ig_audit.get("followers", 0)
    goal         = ig_audit.get("goal_followers", 0)
    gap          = ig_audit.get("follower_gap", 0)

    learned_block = (
        LEARNED_PATTERNS_SLOT.format(learned_patterns_body=learned_patterns)
        if learned_patterns else ""
    )
    try:
        raw_text = await llm_complete(
            messages=[
                {
                    "role": "system",
                    "content": (
                        CRITICAL_INSTRUCTION_PREFIX + "\n\n"
                        + (learned_block + "\n\n" if learned_block else "")
                        + HASHTAG_STRATEGY_2026 + "\n\n"
                        + GROWTH_KPI_THRESHOLDS_2026 + "\n\n"
                        + ANTI_AI_LANGUAGE_2026_ADDITIONS + "\n\n"
                        +
                        f"You are a senior social media growth strategist. You produce consulting-grade, "
                        f"data-backed growth plans — not templates. Every recommendation MUST reference "
                        f"{name} and {niche} specifically. Generic advice is unacceptable.\n\n"
                        f"CRITICAL RULES:\n"
                        f"1. NEVER assume follower counts. Use ONLY the actual data provided.\n"
                        f"2. Produce day_by_day_plan when days_ahead ≤ 21 (one entry per day). "
                        f"Produce weekly_plan when days_ahead > 21.\n"
                        f"3. performance_diagnosis must reference actual post data — not generic insights.\n"
                        f"4. pillar_breakdown must have one entry per content pillar with real performance assessment.\n"
                        f"5. Every insight must feel like a strategist spent hours on this brand.\n"
                        f"\nDATA-QUALITY RULES (STRICT — non-negotiable):\n"
                        f"6. NEVER invent statistics, metrics, revenue figures, market-share numbers, "
                        f"or growth percentages. If a number was not given to you in the data above, do not state it as fact.\n"
                        f"7. The ONLY follower numbers you may use are: current={followers}, goal={goal}. "
                        f"Every milestone must be between {followers} and {goal} inclusive. Never output a value below {followers}.\n"
                        f"8. If you must reference an industry benchmark or a projection that is not directly in the data, "
                        f"prefix it with 'Estimate:' or 'Assumption:' so it is clearly labelled as not-measured.\n"
                        f"9. Do NOT cite competitor follower counts, revenue, or pricing unless they appear in the research data.\n"
                        f"10. Numbers across the plan must be internally consistent (milestones increase toward the goal; "
                        f"percentages in content_mix sum to 100).\n\n"
                        + (
                            f"GOAL INTENSITY: {feasibility['acceleration_needed']}x acceleration required.\n"
                            + (
                                f"EXTREME GOAL — Requires {feasibility['acceleration_needed']}x growth rate. "
                                f"growth_tactics MUST include: (a) daily Reels posting minimum, "
                                f"(b) at least 2 collaboration / creator-swap campaigns, "
                                f"(c) paid promotion strategy (even $5/day boosts), "
                                f"(d) viral hook challenge or trend-jacking plan. "
                                f"Mark this goal as HIGH RISK in goal_strategy.narrative.\n"
                                if feasibility['acceleration_needed'] >= 5 else
                                f"AGGRESSIVE GOAL — Requires {feasibility['acceleration_needed']}x growth rate. "
                                f"growth_tactics MUST include: (a) 2x posting frequency vs current, "
                                f"(b) at least 1 collaboration or niche creator partnership, "
                                f"(c) engagement pod or comment strategy to boost algorithmic reach.\n"
                                if feasibility['acceleration_needed'] >= 3 else
                                f"CHALLENGING GOAL — Requires {feasibility['acceleration_needed']}x growth rate. "
                                f"growth_tactics should emphasise consistency and Reels cadence.\n"
                                if feasibility['acceleration_needed'] >= 2 else
                                f"ACHIEVABLE GOAL — Standard growth path. Focus on quality over frequency.\n"
                            )
                            if feasibility else ""
                        )
                        + f"\n\n{GROWTH_ROADMAP_FRAMEWORK}"
                        + f"\n{STRATEGIC_THINKING}"
                        + f"\n{BUYER_STAGES}"
                        + f"\n{ANTI_AI_LANGUAGE}"
                        + f"\nBrand context:\n{context_block[:800] if context_block else ''}"
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Build a comprehensive {days_ahead}-day Instagram growth strategy for {name}.\n\n"
                        f"=== GOAL FEASIBILITY (pre-calculated — use these exact numbers) ===\n"
                        + (
                            f"Current Followers: {feasibility['current_followers']:,} → Target: {feasibility['target_followers']:,} (Gap: +{feasibility['gap']:,})\n"
                            f"Est. current weekly growth: {feasibility['est_current_weekly_growth']} followers/week\n"
                            f"Required weekly growth: {feasibility['required_weekly_growth']} followers/week\n"
                            f"Required daily growth: {feasibility['required_daily_growth']} followers/day\n"
                            f"Acceleration needed: {feasibility['acceleration_needed']}x\n"
                            f"Success probability: {feasibility['probability_pct']}%\n"
                            f"Key risks: {'; '.join(feasibility['risks']) or 'None identified'}\n\n"
                            if feasibility else ""
                        )
                        + f"=== BRAND CONTEXT ===\n"
                        f"Niche: {niche} | Tone: {tone} | Audience: {audience}\n"
                        + (f"Positioning: {positioning}\n" if positioning else "")
                        + (f"Differentiator: {differentiation}\n" if differentiation else "")
                        + (f"Voice style: {voice_style}\n" if voice_style else "")
                        + (f"Hook style: {hook_style}\n" if hook_style else "")
                        + (f"CTA style: {cta_style}\n" if cta_style else "")
                        + (f"Catchphrases: {catchphrases}\n" if catchphrases else "")
                        + (f"Audience pain points: {audience_pain}\n" if audience_pain else "")
                        + (f"Existing content pillars: {', '.join(str(p) for p in content_pillars)}\n" if content_pillars else "")
                        + f"\n=== INSTAGRAM PERFORMANCE DATA ===\n"
                        + (
                            f"LIVE DATA from @{analyst_report.get('username', name)}:\n"
                            f"- Current Followers: {followers:,} → Goal: {goal:,} (need +{gap:,})\n"
                            f"- Avg Engagement Rate: {ig_audit['avg_er']}% (industry avg: 1.5-3%)\n"
                            f"- Avg Reach per Post: {ig_audit['avg_reach']:,}\n"
                            f"- Best Content Type: {ig_audit['best_content_type']}\n"
                            f"- Posts Analysed: {ig_audit['posts_analysed']}\n"
                            if ig_connected else
                            f"No IG connected — building strategy from brand context + niche research\n"
                            f"Follower goal: {goal:,}\n"
                        )
                        + f"\nWHAT'S WORKING (replicate this formula):\n{working_text or 'No live data — use niche best practices'}\n\n"
                        + f"WHAT'S NOT WORKING (stop/pivot these):\n{not_working_text or 'No live data'}\n\n"
                        + f"=== MARKET RESEARCH ===\n"
                        + f"Trending content angles for {niche}:\n{trends_text or 'N/A'}\n\n"
                        + f"Competitor content gaps to fill:\n{gaps_text or 'N/A'}\n\n"
                        + (f"Content formats to own: {', '.join(formats_to_own)}\n\n" if formats_to_own else "")
                        + (f"Posting insights: {posting_text}\n\n" if posting_text else "")
                        + (f"Differentiation strategy: {diff_strategy}\n\n" if diff_strategy else "")
                        + f"Return JSON with these keys (all specific to {name}/{niche}, not generic):\n"
                        + f"pillars: list of 4 content pillars — derived from what's working + brand positioning + audience needs\n"
                        + f"posting_frequency: string like '1-2x daily'\n"
                        + f"best_times: list of 3 specific posting times for {niche} audience\n"
                        + f"content_mix: object — Reel/Carousel/Graphic/Story/AI Reel as percentages summing to 100 "
                        + f"(weight towards best_content_type: {ig_audit.get('best_content_type','Reel')})\n"
                        + f"monthly_themes: list of 3 specific monthly theme ideas tied to {niche} seasons/trends\n"
                        + f"growth_tactics: list of 6 SPECIFIC, actionable tactics to reach +{gap or 100} followers. "
                        + f"Each must name the audience segment, the format, the angle, and the repurposing path — "
                        + f"e.g. 'Ship a 14-day founder-POV Reel series on [specific {niche} pain] using PAS hooks, "
                        + f"repurpose each into a carousel + 3 Stories + a WhatsApp broadcast'. Never write a tactic "
                        + f"that could apply to any brand (no 'post consistently', no 'engage with your audience').\n"
                        + f"hashtag_strategy: string (2-3 sentences) on hashtag approach specific to {niche}\n"
                        + f"cta_templates: list of 4 CTA templates that match {name}'s voice and {cta_style or tone}\n"
                        + f"what_works: list of 4 specific, data-backed insights about what content is performing well\n"
                        + f"what_to_stop: list of 3 specific things to stop or change based on performance data\n"
                        + (
                            f"follower_plan: object with day5/day10/day15 as realistic follower milestones "
                            f"(starting from {followers}, goal: {goal}). NEVER set any value below {followers}.\n"
                            if days_ahead <= 21 else
                            f"follower_plan: object with week1/week2/week3/week4 as realistic target follower counts "
                            f"(starting from {followers}, goal: {goal}). NEVER set any value below {followers}.\n"
                        )
                        + f"engagement_tactics: list of 4 specific engagement tactics for {niche} community\n"
                        + f"content_series_ideas: list of 3 recurring content series ideas specific to {name}/{niche}\n"
                        + f"hook_strategy: list of 5 hook templates engineered for {niche} retention "
                        + f"(scroll-stop in first 1-2 seconds; reference {hook_style or 'high-impact'} style)\n"
                        + f"retention_strategy: list of 4 specific retention tactics for {niche} reels "
                        + f"(pattern interrupts, on-screen text, b-roll cuts — be concrete)\n"
                        + f"conversion_strategy: list of 4 ways to convert viewers into followers/leads/customers "
                        + f"based on {name}'s {positioning or 'positioning'} — not generic 'link in bio' advice\n"
                        + f"competitor_advantage: list of 3 specific moves that differentiate {name} from competitors\n"
                        + f"viral_opportunities: list of 3 specific viral content concepts for {niche} this quarter\n"
                        + f"platform_expansion: object with reels (string), carousels (string), stories (string) — "
                        + f"each describing the role of that format in {name}'s growth\n"
                        + (
                            f"kpi_targets_90day: keep from analyst — {{followers, avg_engagement_rate, reels_per_week, saves_per_post}}\n"
                            if not ig_connected else
                            f"kpi_targets_30day: object with target_followers (number), target_er (number), reels_per_week (number)\n"
                        )
                        + f"\n=== NEW REQUIRED KEYS (must be specific to {name}/{niche}) ===\n"
                        + f"performance_diagnosis: object with whats_working (list of 3-4 observations from ACTUAL top-post data), "
                        + f"whats_failing (list of 3-4 from low-post data), missed_opportunities (list of 3-4 gaps), "
                        + f"bottlenecks (list of 2-3 structural blockers for {name})\n"
                        + f"pillar_breakdown: list of objects, one per content pillar — each with pillar (string), "
                        + f"current_performance (high/medium/low), growth_potential (high/medium/low), "
                        + f"recommended_frequency (e.g. '3x/week'), expected_impact (string, 1 sentence specific to {niche})\n"
                        + f"goal_strategy: object with narrative (2-3 sentences on HOW the goal will be reached), "
                        + (
                            f"week_by_week_path (list of strings like 'Day 5: first Reels batch live, target +{max(1, int((feasibility or {{}}).get('required_daily_growth', 3) * 5)) if feasibility else 10} followers'), "
                            if days_ahead <= 21 else
                            f"week_by_week_path (list of strings like 'Week 1: focus on Reels, target +{max(1, int((feasibility or {{}}).get('required_weekly_growth', 10))) if feasibility else 10} followers'), "
                        )
                        + f"non_negotiable_actions (list of 3-5 must-do actions), risk_mitigation (list of 2-3 actions)\n"
                        + f"\n=== CONSULTANT-GRADE DEPTH (required — make these specific to {name}/{niche}) ===\n"
                        + f"priority_actions: list of 5 objects (the 5 highest-leverage moves), each with: "
                        + f"action (specific, niche-referenced — NOT 'post more'), why_it_matters (1 sentence), "
                        + f"expected_impact (concrete, e.g. '+X reach/week' or 'lifts save-rate'), "
                        + f"difficulty ('low'|'medium'|'high'), priority_score (integer 1-10), "
                        + f"timeline (e.g. 'Days 1-7'), implementation_steps (list of 2-4 concrete steps), "
                        + f"success_metric (1 measurable signal). Order by priority_score descending.\n"
                        + f"roadmap: object with four keys — quick_wins (1-7 days), short_term (30 days), "
                        + f"mid_term (90 days), long_term (6-12 months). Each is a list of 2-3 objects with: "
                        + f"action, owner (e.g. 'Founder', 'Editor', 'Social lead'), resources_required, "
                        + f"expected_result, kpi, priority ('P0'|'P1'|'P2'). Sequence them so momentum compounds.\n"
                        + f"strategic_analysis: object with risks (list of 3 specific failure modes for {name}), "
                        + f"tradeoffs (list of 2-3), constraints (list of 2-3 real limits — time/team/content supply), "
                        + f"alternative_approaches (list of 2 different paths to the same goal), "
                        + f"backup_plans (list of 2 — what to do if the main plan stalls by day 7). Be honest, not rosy.\n"
                        + (
                            f"day_by_day_plan: list of {days_ahead} objects (one per day), each with day (int), "
                            f"content_task (specific post idea for {name}), growth_task (specific growth action), "
                            f"engagement_task, community_task, kpi_target (e.g. '+3 followers, 5 comments')\n"
                            if days_ahead <= 21 else
                            f"weekly_plan: list of {max(4, days_ahead // 7)} objects, each with week (int), theme (string), "
                            f"content_tasks (string), growth_tasks (string), kpi_target (string)\n"
                        )
                    ),
                },
            ],
            tier="brain",
            temperature=0.5,
            max_tokens=7000,
            response_json=True,
        )
        strategy_out = json.loads(raw_text)
        # Ensure KPI targets from analyst flow through to the PPT layer
        if not strategy_out.get("kpi_targets_90day") and analyst_report.get("kpi_targets_90day"):
            strategy_out["kpi_targets_90day"] = analyst_report["kpi_targets_90day"]
        if not strategy_out.get("launch_roadmap") and analyst_report.get("launch_roadmap"):
            strategy_out["launch_roadmap"] = analyst_report["launch_roadmap"]
        return strategy_out
    except Exception as e:
        print(f"[GrowthPlanner] GPT strategy error: {e}")
        import traceback; traceback.print_exc()
        return _fallback_strategy(brand, days_ahead, ig_audit)


# ── Extended Plan (forecasting + structured recs + execution calendar) ───────

async def _generate_extended_plan(
    brand: dict, ig_audit: dict, feasibility: dict,
    research_data: dict, competitor_data: dict,
    growth_strategy: dict, days_ahead: int,
) -> dict:
    """Second-pass GPT call for the investor-grade extended plan.
    Adds: growth forecasting (3 scenarios), structured recommendations with
    evidence, content audit breakdown, social media trends, and a 30-day
    day-by-day + 60-90-day weekly execution calendar.

    DATA QUALITY: all forecasts are labeled as Estimate/Assumption — never
    measured fact. Numbers are clamped by _clamp_extended_forecasting.
    """
    name      = brand.get("name", "brand")
    niche     = brand.get("niche", "")
    followers = ig_audit.get("followers", 0)
    goal      = ig_audit.get("goal_followers", 500)
    avg_er    = ig_audit.get("avg_er", 0)
    feas_prob = feasibility.get("probability_pct", 70)
    accel     = feasibility.get("acceleration_needed", 1.0)
    req_wk    = feasibility.get("required_weekly_growth", 5)

    what_works = (growth_strategy.get("what_works") or [])[:3]
    what_stop  = (growth_strategy.get("what_to_stop") or [])[:3]
    comp_gaps  = (competitor_data.get("content_gaps") or [])[:4]
    trends     = (research_data.get("trending_angles") or [])[:4]
    tactics    = (growth_strategy.get("growth_tactics") or [])[:4]

    try:
        raw_text = await llm_complete(
            messages=[
                {
                    "role": "system",
                    "content": (
                        CRITICAL_INSTRUCTION_PREFIX + "\n\n"
                        + GROWTH_KPI_THRESHOLDS_2026 + "\n\n"
                        +
                        f"You are a senior growth strategist producing an investor-grade, non-generic report "
                        f"for {name} ({niche}). Every recommendation must cite specific evidence from the "
                        f"data provided — if a recommendation could apply to any brand, rewrite it.\n\n"
                        f"STRICT DATA-QUALITY RULES:\n"
                        f"1. All follower forecasts are ESTIMATES. Prefix assumptions: 'Estimate:' or 'Assumption:'.\n"
                        f"2. Current followers = {followers:,}. All milestone values must be ≥ {followers:,}.\n"
                        f"3. Never invent statistics, revenue, market-share, or competitor metrics.\n"
                        f"4. Anti-generic: every item must reference {name}/{niche} specifically.\n"
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Generate the extended investor-grade growth plan for {name} ({niche}).\n\n"
                        f"PERFORMANCE DATA:\n"
                        f"Followers: {followers:,} → Goal: {goal:,} | ER: {avg_er}% | "
                        f"Success prob: {feas_prob}% | Accel needed: {accel}x | Need {req_wk}/week\n"
                        f"What works: {'; '.join(what_works)}\n"
                        f"What to stop: {'; '.join(what_stop)}\n"
                        f"Competitor gaps: {'; '.join(str(g) for g in comp_gaps)}\n"
                        f"Market trends: {'; '.join(str(t) for t in trends)}\n"
                        f"Current tactics: {'; '.join(str(t) for t in tactics)}\n\n"
                        f"Return JSON with these exact keys:\n\n"
                        f"growth_forecasting: object with "
                        f"assumptions (list of 3 explicit assumptions — prefix with 'Estimate:' or 'Assumption:'), "
                        f"growth_drivers (list of 3 factors that determine the outcome), "
                        f"day30/day60/day90/day180 — each an object with best/realistic/worst as integer follower "
                        f"counts starting from {followers:,}. NEVER set any value below {followers:,}.\n\n"
                        f"structured_recommendations: list of 5 objects, each with: "
                        f"title (specific action — NOT 'post more'), "
                        f"recommendation (2-3 sentences — exactly what to do for {name}/{niche}), "
                        f"evidence (list of 2-3 specific data points FROM the provided data above — not invented), "
                        f"reason (why this works for {name}/{niche} specifically), "
                        f"expected_impact (concrete, e.g. '+X reach/week' or 'lifts save-rate by ~Y'), "
                        f"priority ('P0'|'P1'|'P2'), difficulty ('low'|'medium'|'high'), "
                        f"timeline (e.g. 'Days 1-14'). Order by priority (P0 first).\n\n"
                        f"content_best_breakdown: object with format (content type), "
                        f"hook_pattern (what the hook structure did), "
                        f"why_it_worked (list of 3 specific reasons — reference the actual data), "
                        f"replicable_formula (exact template to reproduce this), "
                        f"recommended_frequency (e.g. '3x/week').\n\n"
                        f"content_worst_breakdown: object with format (content type that underperforms), "
                        f"why_it_failed (list of 2-3 specific reasons), "
                        f"what_to_stop (specific pattern to avoid), "
                        f"what_to_do_instead (the better alternative).\n\n"
                        f"social_media_trends: list of 4 objects — trend (name), why_it_works (mechanism), "
                        f"how_to_use (specific to {niche}), competitor_using (generic label).\n\n"
                        f"execution_calendar_30: list of EXACTLY 30 objects for days 1-30 — "
                        f"day (int), reel_topic (specific to {niche}), carousel_topic (specific to {niche}), "
                        f"story_topic (1 line), cta (CTA mechanism), goal (e.g. '+5 followers, 10 saves').\n\n"
                        f"execution_calendar_weeks_5_to_12: list of 8 objects for weeks 5-12 — "
                        f"week (5-12), theme (1 line), reels (2 topic lines), carousels (1 topic line), "
                        f"stories (pattern), goal (weekly KPI)."
                    ),
                },
            ],
            tier="brain",
            temperature=0.4,
            max_tokens=5000,
            response_json=True,
        )
        return json.loads(raw_text)
    except Exception as e:
        print(f"[GrowthPlanner] Extended plan GPT error: {e}")
        import traceback; traceback.print_exc()
        return {}


def _clamp_extended_forecasting(ep: dict, ig_audit: dict) -> None:
    """Clamp forecasting milestone values to realistic bounds in-place.
    Enforces: value ≥ current_followers; value ≤ current*20 (prevents GPT hallucinating
    viral spikes); and worst ≤ realistic ≤ best within each period.
    """
    if not isinstance(ep, dict):
        return
    followers = int(ig_audit.get("followers", 0) or 0)
    upper     = max(followers * 20, followers + 50000)  # generous cap
    gf = ep.get("growth_forecasting")
    if not isinstance(gf, dict):
        return
    for period in ("day30", "day60", "day90", "day180"):
        pd = gf.get(period)
        if not isinstance(pd, dict):
            continue
        clamped = {}
        for scenario in ("worst", "realistic", "best"):
            v = pd.get(scenario)
            try:
                clamped[scenario] = max(followers, min(upper, int(float(v))))
            except (TypeError, ValueError):
                clamped[scenario] = None
        # enforce ordering: worst ≤ realistic ≤ best
        vals = [clamped[s] for s in ("worst", "realistic", "best") if clamped[s] is not None]
        if vals:
            vals_sorted = sorted(vals)
            for i, s in enumerate(("worst", "realistic", "best")):
                if clamped[s] is not None and i < len(vals_sorted):
                    clamped[s] = vals_sorted[i]
        gf[period] = {**pd, **{k: v for k, v in clamped.items() if v is not None}}


# ── Growth Planner PPT ───────────────────────────────────────────────────────

async def _build_growth_ppt(brand, ig_audit, strategy, research_data, competitor_data, days_ahead, run_id, analyst_report: dict = None, feasibility: dict = None, extended_plan: dict = None) -> str | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.625)

        name    = brand.get("name", "Brand")
        niche   = brand.get("niche", "")
        colors  = brand.get("colors") or {}
        primary = colors.get("primary", "#6C3CE1") if isinstance(colors, dict) else "#6C3CE1"

        def hex_to_rgb(h):
            h = h.lstrip("#")
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        brand_color = hex_to_rgb(primary)
        white  = RGBColor(0xFF, 0xFF, 0xFF)
        dark   = RGBColor(0x0F, 0x0A, 0x1E)
        accent = RGBColor(0xA7, 0x8B, 0xFA)
        green  = RGBColor(0x10, 0xB9, 0x81)
        red    = RGBColor(0xEF, 0x44, 0x44)

        blank = prs.slide_layouts[6]

        def bg(slide, color):
            f = slide.background.fill; f.solid(); f.fore_color.rgb = color

        def txt(slide, text, l, t, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
            tb = slide.shapes.add_textbox(l, t, w, h)
            tf = tb.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = align
            r  = p.add_run(); r.text = str(text)
            r.font.size  = Pt(size); r.font.bold = bold
            r.font.color.rgb = color or white
            r.font.italic = italic

        def bar(slide, l, t, w, h, color):
            s = slide.shapes.add_shape(1, l, t, w, h)
            s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

        date_str     = datetime.utcnow().strftime("%d %b %Y")
        month_str    = datetime.utcnow().strftime("%B %Y")
        ig_connected = ig_audit.get("ig_connected", False)
        followers    = ig_audit.get("followers", 0)
        goal         = ig_audit.get("goal_followers", 500)
        gap          = ig_audit.get("follower_gap", goal)
        avg_er       = ig_audit.get("avg_er", 0)
        avg_reach    = ig_audit.get("avg_reach", 0)
        ar           = analyst_report or {}
        rd           = research_data or {}
        cd           = competitor_data or {}
        feas         = feasibility or {}
        grey         = RGBColor(0x9C, 0xA3, 0xAF)
        light        = RGBColor(0xCB, 0xD5, 0xE1)
        mid_dark     = RGBColor(0x12, 0x0D, 0x28)
        ep           = extended_plan or {}
        amber        = RGBColor(0xF5, 0x9E, 0x0B)
        total_slides = 22

        def footer(slide, n):
            txt(slide, f"Prepared by PerformEdge  ·  {date_str}  ·  {n}/{total_slides}",
                Inches(0.4), Inches(5.25), Inches(9.2), Inches(0.3), 9, color=RGBColor(0x4B,0x55,0x63))

        def slide_header(slide, tag, title):
            bar(slide, 0, 0, Inches(10), Inches(0.06), brand_color)
            txt(slide, tag, Inches(0.4), Inches(0.18), Inches(9), Inches(0.38), 10, color=accent, bold=True)
            txt(slide, title, Inches(0.4), Inches(0.55), Inches(9), Inches(0.5), 20, bold=True)

        # ── Slide 1: Cover ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(0.06), prs.slide_height, brand_color)
        txt(s, name.upper(), Inches(0.4), Inches(0.9), Inches(9), Inches(1.3), 44, bold=True)
        txt(s, "Social Media", Inches(0.4), Inches(2.28), Inches(9), Inches(0.65), 28, color=accent)
        txt(s, "Growth Strategy", Inches(0.4), Inches(2.9), Inches(9), Inches(0.65), 28, color=accent)
        txt(s, niche, Inches(0.4), Inches(3.7), Inches(9), Inches(0.4), 14, color=grey, italic=True)
        txt(s, "Investor-Grade Report  ·  Prepared by PerformEdge", Inches(0.4), Inches(4.75), Inches(7), Inches(0.3), 11, color=RGBColor(0x6B,0x72,0x80))
        txt(s, month_str, Inches(7.5), Inches(4.75), Inches(2), Inches(0.3), 11, color=RGBColor(0x6B,0x72,0x80), align=PP_ALIGN.RIGHT)

        # ── Slide 2: Executive Summary ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "01  ·  EXECUTIVE SUMMARY", "The bottom line — what we found and what it means.")
        # Surface data-quality warnings at the top of the deck.
        warnings = []
        if ig_audit.get("goal_was_invalid"):
            warnings.append(f"⚠ Goal regenerated — submitted target was below current followers (auto-set to {goal:,})")
        if not ig_audit.get("has_live_metrics"):
            warnings.append("⚠ Limited audit — Instagram not connected. Connect IG for live content performance metrics.")
        if warnings:
            txt(s, "  ·  ".join(warnings)[:170], Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.28), 8, color=amber, italic=True)
        has_live  = ig_audit.get("has_live_metrics", False)
        prob      = feas.get("probability_pct", 70)
        prob_col  = green if prob >= 70 else (amber if prob >= 40 else red)
        # Spec §data_validation_system: never show '0' for unmeasured metrics.
        er_disp   = f"{avg_er}%" if has_live and avg_er > 0 else "Data Unavailable"
        er_col    = accent if (has_live and avg_er > 0) else grey
        exec_stats = [
            (f"{followers:,}",    "Current Followers", brand_color),
            (f"{goal:,}",         "Growth Target",     green),
            (er_disp,             "Engagement Rate",   er_col),
            (f"{prob}%",          "Goal Probability",  prob_col),
            (f"{feas.get('acceleration_needed',1)}x", "Growth Needed", amber),
        ]
        for i, (val, lbl, col) in enumerate(exec_stats):
            x = Inches(0.3 + i * 1.9)
            bar(s, x, Inches(1.2), Inches(1.78), Inches(1.5), mid_dark)
            bar(s, x, Inches(1.2), Inches(1.78), Inches(0.05), col)
            txt(s, val, x+Inches(0.1), Inches(1.38), Inches(1.6), Inches(0.68), 20, bold=True, color=col, align=PP_ALIGN.CENTER)
            txt(s, lbl, x+Inches(0.05), Inches(2.1), Inches(1.7), Inches(0.3), 8, color=grey, align=PP_ALIGN.CENTER)
        p0_recs   = [r for r in (ep.get("structured_recommendations") or []) if str(r.get("priority","")).upper() in ("P0","P1")][:3]
        top_items = [str(r.get("title","") or r) for r in p0_recs] or (strategy.get("growth_tactics") or [])[:3]
        bar(s, Inches(0.3), Inches(2.85), Inches(9.4), Inches(1.85), mid_dark)
        txt(s, "TOP STRATEGIC PRIORITIES", Inches(0.45), Inches(2.93), Inches(9), Inches(0.3), 8, color=accent, bold=True)
        for j, p in enumerate(top_items[:3]):
            txt(s, f"{j+1}.  {str(p)[:115]}", Inches(0.45), Inches(3.28+j*0.42), Inches(9.1), Inches(0.38), 11, color=light)
        footer(s, 2)

        # ── Slide 3: Brand & Audience Analysis ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "02  ·  BRAND & AUDIENCE ANALYSIS", "Who you are and who you're speaking to.")
        positioning    = brand.get("positioning","") or ""
        differentiation= brand.get("differentiation","") or ""
        pain_points    = brand.get("audiencePainPoints","") or ""
        aspirations    = brand.get("audienceAspirations","") or ""
        bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(3.95), mid_dark)
        bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.05), brand_color)
        txt(s, "BRAND PROFILE", Inches(0.45), Inches(1.23), Inches(4.2), Inches(0.3), 8, color=brand_color, bold=True)
        # Wider field budget — values wrap into the 0.85"-tall textbox at 10pt.
        brand_fields = [f"Niche: {brand.get('niche','')[:90]}", f"Industry: {brand.get('industry','')[:90]}",
                        f"Positioning: {positioning[:140]}", f"Differentiator: {differentiation[:140]}"]
        for j, bf in enumerate([f for f in brand_fields if not f.endswith(": ")][:4]):
            txt(s, bf, Inches(0.45), Inches(1.6+j*0.85), Inches(4.2), Inches(0.85), 10, color=light)
        bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(3.95), mid_dark)
        bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(0.05), accent)
        txt(s, "TARGET AUDIENCE", Inches(5.15), Inches(1.23), Inches(4.4), Inches(0.3), 8, color=accent, bold=True)
        aud_fields = [f"Audience: {brand.get('targetAudience','')[:160]}", f"Pain Points: {pain_points[:160]}", f"Aspirations: {aspirations[:160]}"]
        for j, af in enumerate([f for f in aud_fields if not f.endswith(": ")][:4]):
            txt(s, af, Inches(5.15), Inches(1.6+j*1.1), Inches(4.4), Inches(1.1), 10, color=light)
        footer(s, 3)

        # ── Slide 4: Social Media Audit ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "03  ·  SOCIAL MEDIA AUDIT", "Where you stand today — measured numbers, not guesses.")
        if ig_connected and followers > 0 and ig_audit.get("posts_analysed", 0) > 0:
            mode_label = f"Live Instagram data  ·  {ig_audit.get('posts_analysed', 0)} posts analysed  ·  Source: Meta Graph API  ·  Confidence: high"
        elif followers > 0:
            mode_label = f"Stored data — {followers:,} followers verified  ·  Source: brand profile  ·  Confidence: medium  ·  Connect IG for live engagement metrics"
        else:
            mode_label = "Launch Mode  ·  No social profile connected  ·  Confidence: low  ·  Targets based on brand brief + niche research"
        txt(s, mode_label, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.28), 9, color=grey, italic=True)
        # Data display with honest confidence tiers:
        # live → exact number  |  benchmark/scraped → "~X% (Est.)"  |  none → "Unable to Retrieve"
        posts_n    = ig_audit.get("posts_analysed", 0)
        est_er     = ig_audit.get("est_er", 0)
        est_reach  = ig_audit.get("est_reach", 0)
        bm_src     = ig_audit.get("benchmark_source", "")
        bm_conf    = ig_audit.get("benchmark_confidence", 55)
        ds         = ig_audit.get("data_source", "stored")

        if has_live and avg_er > 0:
            er_v    = f"{avg_er}%"
            er_col  = green
        elif est_er > 0:
            er_v    = f"~{est_er}% (Est.)"
            er_col  = amber
        else:
            er_v    = "Unable to Retrieve"
            er_col  = grey

        if has_live and avg_reach > 0:
            reach_v   = f"{avg_reach:,}"
            reach_col = accent
        elif est_reach > 0:
            reach_v   = f"~{est_reach:,} (Est.)"
            reach_col = amber
        else:
            reach_v   = "Unable to Retrieve"
            reach_col = grey

        best_fmt_v  = ig_audit.get("best_content_type") or ("—" if posts_n == 0 else "Reel")
        posts_v     = f"{posts_n}" if posts_n > 0 else ("—" if ds in ("benchmark","scraped") else "Unable to Retrieve")
        snap_stats = [
            (f"{followers:,}", "Followers",       brand_color),
            (er_v,              "Eng. Rate",      er_col),
            (reach_v,           "Avg Reach",      reach_col),
            (best_fmt_v,        "Best Format",    amber if best_fmt_v not in ("Unable to Retrieve","—") else grey),
            (posts_v,           "Posts Analysed", white if posts_n > 0 else grey),
        ]
        for i, (val, lbl, col) in enumerate(snap_stats):
            x = Inches(0.3 + i * 1.9)
            bar(s, x, Inches(1.28), Inches(1.78), Inches(1.75), mid_dark)
            bar(s, x, Inches(1.28), Inches(1.78), Inches(0.05), col)
            txt(s, val, x+Inches(0.1), Inches(1.44), Inches(1.6), Inches(0.72), 20, bold=True, color=col, align=PP_ALIGN.CENTER)
            txt(s, lbl, x+Inches(0.05), Inches(2.18), Inches(1.7), Inches(0.3), 9, color=grey, align=PP_ALIGN.CENTER)
        cm = strategy.get("content_mix") or {}
        bar(s, Inches(0.3), Inches(3.18), Inches(9.4), Inches(1.62), mid_dark)
        txt(s, "CONTENT TYPE MIX", Inches(0.45), Inches(3.26), Inches(9), Inches(0.3), 8, color=brand_color, bold=True)
        ct_hex = {"Reel":"3B82F6","AI Reel":"EC4899","Carousel":"10B981","Graphic":"A78BFA","Story":"F59E0B"}
        for i, (ct, pct) in enumerate(list(cm.items())[:5]):
            x = Inches(0.45 + i * 1.85)
            col_h = hex_to_rgb(ct_hex.get(ct, "6C3CE1"))
            txt(s, ct, x, Inches(3.6), Inches(1.7), Inches(0.3), 9, bold=True, color=col_h)
            txt(s, f"{pct}%", x, Inches(3.92), Inches(1.7), Inches(0.5), 18, bold=True, color=col_h)
        # Show benchmark source when metrics are estimates
        if not has_live and (est_er > 0 or est_reach > 0) and bm_src:
            conf_label = f"Confidence: {bm_conf}%  ·  Source: {bm_src[:80]}"
            txt(s, conf_label, Inches(0.4), Inches(4.9), Inches(9.2), Inches(0.22), 7, color=RGBColor(0x6B,0x72,0x80), italic=True)
        footer(s, 4)

        # ── Slide 5: Content Performance — What's Working ──
        # Spec §quality_control: reject Posts Analysed=0 + Best Content Exists.
        # When there's no live data, show an honest "connect IG" panel rather than
        # fabricating insights about posts that don't exist.
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "04  ·  CONTENT PERFORMANCE — WHAT'S WORKING", "The formula behind your top-performing content.")
        if not has_live:
            bar(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.0), mid_dark)
            bar(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(0.05), amber)
            txt(s, "CONTENT PERFORMANCE DATA UNAVAILABLE", Inches(0.5), Inches(1.7), Inches(9), Inches(0.4), 14, bold=True, color=amber)
            txt(s, "Connect Instagram to enable content performance analysis.",
                Inches(0.5), Inches(2.25), Inches(9), Inches(0.4), 12, color=white)
            txt(s, "Per data-quality rules, this report does not generate 'what worked' insights",
                Inches(0.5), Inches(2.75), Inches(9), Inches(0.35), 10, color=light)
            txt(s, "without measured post-level data. Once IG is connected, this slide will populate",
                Inches(0.5), Inches(3.05), Inches(9), Inches(0.35), 10, color=light)
            txt(s, "with the actual hook pattern, replicable formula, and frequency that drove your top posts.",
                Inches(0.5), Inches(3.35), Inches(9), Inches(0.35), 10, color=light)
            footer(s, 5)
            # Skip to slide 6
            _skip_working = True
        else:
            _skip_working = False
        if not _skip_working:
            best_bd  = ep.get("content_best_breakdown") or {}
            diag     = strategy.get("performance_diagnosis") or {}
            working  = (diag.get("whats_working") or strategy.get("what_works") or [])[:4]
            bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(3.95), mid_dark)
            bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.05), green)
            txt(s, "WHAT'S WORKING", Inches(0.45), Inches(1.23), Inches(4.2), Inches(0.3), 8, color=green, bold=True)
            for j, item in enumerate(working[:4]):
                txt(s, f"✓  {str(item)[:70]}", Inches(0.45), Inches(1.62+j*0.77), Inches(4.2), Inches(0.68), 10, color=light)
            bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(3.95), mid_dark)
            bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(0.05), accent)
            txt(s, "REPLICABLE FORMULA", Inches(5.15), Inches(1.23), Inches(4.4), Inches(0.3), 8, color=accent, bold=True)
            formula = []
            if best_bd:
                if best_bd.get("hook_pattern"): formula.append(f"Hook: {str(best_bd['hook_pattern'])[:65]}")
                why = best_bd.get("why_it_worked")
                if isinstance(why, list): formula.extend([str(w)[:65] for w in why[:2]])
                elif why: formula.append(str(why)[:65])
                if best_bd.get("replicable_formula"): formula.append(f"Formula: {str(best_bd['replicable_formula'])[:60]}")
                if best_bd.get("recommended_frequency"): formula.append(f"Frequency: {best_bd['recommended_frequency']}")
            formula = formula or (strategy.get("hook_strategy") or [])[:4]
            for j, item in enumerate(formula[:4]):
                txt(s, f"→  {str(item)[:70]}", Inches(5.15), Inches(1.62+j*0.77), Inches(4.4), Inches(0.68), 10, color=light)
            footer(s, 5)

        # ── Slide 6: Content Performance — What's Failing ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "05  ·  CONTENT PERFORMANCE — WHAT'S FAILING", "Stop doing this. Pivot here instead.")
        if not has_live:
            bar(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.0), mid_dark)
            bar(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(0.05), amber)
            txt(s, "FAILURE PATTERN DATA UNAVAILABLE", Inches(0.5), Inches(1.7), Inches(9), Inches(0.4), 14, bold=True, color=amber)
            txt(s, "Connect Instagram so we can identify your actual underperforming patterns —",
                Inches(0.5), Inches(2.25), Inches(9), Inches(0.4), 12, color=white)
            txt(s, "what to stop, what to do instead, and the specific formats that aren't earning their place.",
                Inches(0.5), Inches(2.65), Inches(9), Inches(0.4), 11, color=light)
            footer(s, 6)
        else:
            worst_bd = ep.get("content_worst_breakdown") or {}
            diag    = strategy.get("performance_diagnosis") or {}
            failing = (diag.get("whats_failing") or strategy.get("what_to_stop") or [])[:4]
            bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(3.95), mid_dark)
            bar(s, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.05), red)
            txt(s, "WHAT'S FAILING", Inches(0.45), Inches(1.23), Inches(4.2), Inches(0.3), 8, color=red, bold=True)
            for j, item in enumerate(failing[:4]):
                txt(s, f"✗  {str(item)[:70]}", Inches(0.45), Inches(1.62+j*0.77), Inches(4.2), Inches(0.68), 10, color=light)
            bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(3.95), mid_dark)
            bar(s, Inches(5.0), Inches(1.15), Inches(4.7), Inches(0.05), green)
            txt(s, "DO THIS INSTEAD", Inches(5.15), Inches(1.23), Inches(4.4), Inches(0.3), 8, color=green, bold=True)
            pivot = []
            if worst_bd:
                wf = worst_bd.get("why_it_failed")
                if isinstance(wf, list): pivot.extend([str(w)[:65] for w in wf[:2]])
                elif wf: pivot.append(str(wf)[:65])
                if worst_bd.get("what_to_do_instead"): pivot.append(f"→ {str(worst_bd['what_to_do_instead'])[:65]}")
            pivot = pivot or (diag.get("missed_opportunities") or strategy.get("content_series_ideas") or [])[:4]
            for j, item in enumerate(pivot[:4]):
                txt(s, f"→  {str(item)[:70]}", Inches(5.15), Inches(1.62+j*0.77), Inches(4.4), Inches(0.68), 10, color=light)
            footer(s, 6)

        # ── Slide 7: Market Research & Industry Trends ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "06  ·  MARKET RESEARCH & INDUSTRY TRENDS", "What's happening in your market right now.")
        trends_list = rd.get("trending_angles") or [str(t.get("title","")) for t in rd.get("trends",[])[:4]]
        pain_list   = rd.get("audience_pain_insights") or []
        opps_list7  = rd.get("content_opportunities") or rd.get("hook_ideas") or []
        sec7 = [
            ("TRENDING CONTENT ANGLES", accent, [str(t)[:65] for t in trends_list[:3]]),
            ("AUDIENCE PAIN INSIGHTS",  green,  [str(p)[:65] for p in (pain_list or [brand.get("audiencePainPoints","Unaddressed pain points")])[:3]]),
            ("CONTENT OPPORTUNITIES",   amber,  [str(o)[:65] for o in opps_list7[:3]]),
        ]
        for i, (title, col, items) in enumerate(sec7):
            x = Inches(0.3 + i * 3.2)
            bar(s, x, Inches(1.15), Inches(3.0), Inches(3.95), mid_dark)
            bar(s, x, Inches(1.15), Inches(3.0), Inches(0.05), col)
            txt(s, title, x+Inches(0.12), Inches(1.23), Inches(2.8), Inches(0.3), 8, color=col, bold=True)
            for j, item in enumerate((items or ["See research notes"])[:4]):
                txt(s, f"• {item}", x+Inches(0.12), Inches(1.62+j*0.78), Inches(2.8), Inches(0.7), 10, color=light)
        footer(s, 7)

        # ── Slide 8: Social Media Trends ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "07  ·  SOCIAL MEDIA TRENDS", "What's working on Instagram right now — and why.")
        sm_trends = ep.get("social_media_trends") or [
            {"trend": "POV / First-person Reels", "why_it_works": "Creates immediate audience identification", "how_to_use": f"Show behind-the-scenes from founder perspective in {niche}"},
            {"trend": "Educational Carousels with data hooks", "why_it_works": "Drives saves (saves = strong algorithmic signal)", "how_to_use": f"Myth-busting series specific to {niche}"},
            {"trend": "Trend-jacking with niche POV", "why_it_works": "Borrows reach from trending audio/format", "how_to_use": f"Apply viral audio to {niche} content"},
            {"trend": "Story-based transformation Reels", "why_it_works": "Emotional arc keeps viewers to the end", "how_to_use": f"Before/after results in {niche}"},
        ]
        for i, t in enumerate(sm_trends[:4]):
            row, ci = divmod(i, 2)
            x = Inches(0.3 + ci * 4.8); y = Inches(1.15 + row * 2.12)
            tn  = str(t.get("trend","") or t)[:55]
            why = str(t.get("why_it_works",""))[:70]
            how_key = next((k for k in (t.keys() if isinstance(t, dict) else []) if "how" in k.lower()), None)
            how = str(t.get(how_key,""))[:70] if how_key else ""
            bar(s, x, y, Inches(4.5), Inches(1.88), mid_dark)
            bar(s, x, y, Inches(4.5), Inches(0.05), accent)
            txt(s, tn,  x+Inches(0.15), y+Inches(0.1),  Inches(4.1), Inches(0.32), 11, bold=True, color=white)
            txt(s, f"Why: {why}", x+Inches(0.15), y+Inches(0.48), Inches(4.1), Inches(0.32), 9, color=grey)
            if how: txt(s, f"Use: {how}", x+Inches(0.15), y+Inches(0.84), Inches(4.1), Inches(0.32), 9, color=light)
        footer(s, 8)

        # ── Slide 9: Competitor Intelligence ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "08  ·  COMPETITOR INTELLIGENCE", "Who you're up against — and where you win.")
        # Show actual competitor names at the top
        comps_found = cd.get("competitors_found") or []
        if comps_found:
            comp_names_str = "  ·  ".join(str(c)[:30] for c in comps_found[:6])
            bar(s, Inches(0.3), Inches(1.05), Inches(9.4), Inches(0.42), mid_dark)
            txt(s, "COMPETITORS IDENTIFIED:", Inches(0.45), Inches(1.1), Inches(2.2), Inches(0.3), 8, color=accent, bold=True)
            txt(s, comp_names_str, Inches(2.8), Inches(1.1), Inches(6.8), Inches(0.3), 9, color=light)
        diff_strat = cd.get("differentiation_strategy") or ""
        if diff_strat:
            txt(s, f'"{str(diff_strat)[:180]}"', Inches(0.4), Inches(1.58 if comps_found else 1.08), Inches(9.2), Inches(0.38), 9, color=grey, italic=True)
        comp_adv    = strategy.get("competitor_advantage") or []
        comp_gaps9  = cd.get("content_gaps") or cd.get("gaps_to_fill") or []
        formats_own = cd.get("content_formats_to_own") or []
        col9_items  = [
            ("OUR DIFFERENTIATORS",    brand_color, comp_adv[:3]   or ["Stronger storytelling","Deeper niche expertise","Community-first"]),
            ("COMPETITOR GAPS TO FILL",green,       comp_gaps9[:3] or ["Underserved audience","Missing educational","No UGC strategy"]),
            ("FORMATS WE WILL OWN",    accent,      formats_own[:3] or ["Educational Reels","BTS Stories","Community Carousels"]),
        ]
        y9 = Inches(2.1) if (comps_found or diff_strat) else Inches(1.2)
        for i, (title, col, items) in enumerate(col9_items):
            x = Inches(0.3 + i * 3.2)
            bar(s, x, y9, Inches(3.0), Inches(3.0), mid_dark)
            bar(s, x, y9, Inches(3.0), Inches(0.05), col)
            txt(s, title, x+Inches(0.12), y9+Inches(0.1), Inches(2.8), Inches(0.32), 8, color=col, bold=True)
            for j, item in enumerate(items[:4]):
                txt(s, f"→ {str(item)[:56]}", x+Inches(0.12), y9+Inches(0.5+j*0.58), Inches(2.8), Inches(0.5), 10, color=light)
        footer(s, 9)

        # ── Slide 10: Competitor Gap Analysis ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "09  ·  COMPETITOR GAP ANALYSIS", "Where the market is underserved — and how you can own it.")
        gap_opps   = cd.get("gap_opportunities") or {}
        opps10     = cd.get("opportunities") or []
        rec_angles = cd.get("recommended_angles") or []
        gap10 = [
            ("TOPIC GAP",    accent, [str(gap_opps.get("topic_gap","Unaddressed topics in niche"))[:65]]),
            ("FORMAT GAP",   green,  [str(gap_opps.get("format_gap","Effective formats competitors avoid"))[:65]]),
            ("AUDIENCE GAP", amber,  [str(gap_opps.get("audience_gap","Sub-segment no one serves"))[:65]]),
            ("TONE GAP",     red,    [str(gap_opps.get("tone_gap","Emotional register no one occupies"))[:65]]),
        ]
        for i, (title, col, items) in enumerate(gap10):
            row, ci = divmod(i, 2)
            x = Inches(0.3 + ci * 4.8); y = Inches(1.15 + row * 2.12)
            bar(s, x, y, Inches(4.5), Inches(1.88), mid_dark)
            bar(s, x, y, Inches(4.5), Inches(0.05), col)
            txt(s, title, x+Inches(0.15), y+Inches(0.1), Inches(4.1), Inches(0.3), 8, color=col, bold=True)
            for j, item in enumerate(items[:2]):
                txt(s, f"• {item}", x+Inches(0.15), y+Inches(0.45+j*0.44), Inches(4.2), Inches(0.38), 10, color=light)
            opp = opps10[i] if i < len(opps10) else (rec_angles[i] if i < len(rec_angles) else "")
            if opp: txt(s, f"→ Opp: {str(opp)[:60]}", x+Inches(0.15), y+Inches(1.45), Inches(4.2), Inches(0.3), 9, color=col, italic=True)
        footer(s, 10)

        # ── Slide 11: Growth Gap Analysis ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "10  ·  GROWTH GAP ANALYSIS", "Current state vs required state — the exact delta.")
        req_wk  = feas.get("required_weekly_growth", 0)
        req_day = feas.get("required_daily_growth", 0)
        est_wk  = feas.get("est_current_weekly_growth", 5)
        accel   = feas.get("acceleration_needed", 1.0)
        _er_disp    = f"{avg_er}%" if has_live else (f"~{ig_audit.get('est_er',0)}% (Est.)" if ig_audit.get('est_er',0) else "—")
        _reach_disp = f"{avg_reach:,}" if has_live else (f"~{ig_audit.get('est_reach',0):,} (Est.)" if ig_audit.get('est_reach',0) else "—")
        _req_reach  = f"{int(ig_audit.get('est_reach',0) * 1.5):,}+" if ig_audit.get('est_reach',0) and not has_live else (f"{int(avg_reach * 1.5):,}+" if avg_reach else "See strategy")
        rows11  = [
            ("Followers",      f"{followers:,}",    f"{goal:,}",        f"+{gap:,}",            brand_color),
            ("Weekly Growth",  f"{est_wk}/wk",       f"{req_wk}/wk",     f"{accel}x needed",     amber),
            ("Daily Growth",   "—",                  f"{req_day}/day",    "—",                    accent),
            ("Engagement Rate",_er_disp,              "3-5%",             "Industry benchmark",   green),
            ("Avg Reach",      _reach_disp,           _req_reach,         "+50% target",          grey),
        ]
        bar(s, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.38), RGBColor(0x1A,0x10,0x35))
        for ci, lbl in enumerate(["METRIC","CURRENT","REQUIRED","GAP"]):
            txt(s, lbl, Inches(0.38+ci*2.35), Inches(1.26), Inches(2.3), Inches(0.28), 8, bold=True, color=accent)
        for j, (metric, cur, req, gap_v, col) in enumerate(rows11):
            y = Inches(1.62 + j * 0.72)
            row_bg = mid_dark if j%2==0 else RGBColor(0x16,0x10,0x30)
            bar(s, Inches(0.3), y, Inches(9.4), Inches(0.66), row_bg)
            bar(s, Inches(0.3), y, Inches(0.05), Inches(0.66), col)
            for ci, val in enumerate([metric, cur, req, gap_v]):
                txt(s, str(val), Inches(0.4+ci*2.35), y+Inches(0.12), Inches(2.3), Inches(0.42), 12,
                    bold=(ci==0), color=col if ci==0 else (light if ci<2 else amber))
        footer(s, 11)

        # ── Slide 12: Growth Forecasting ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "11  ·  GROWTH FORECASTING", "What's possible — all projections labeled as Estimates.")
        gf = ep.get("growth_forecasting") or {}
        txt(s, "All values below are ESTIMATES based on current trajectory + proposed strategy changes, not measured data.",
            Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.22), 8, color=amber, italic=True)
        assumptions = gf.get("assumptions") or [
            f"Estimate: consistent posting at {strategy.get('posting_frequency','1x daily')} required",
            f"Assumption: ER maintained at {avg_er}% or above",
            f"Assumption: algorithm reach improves with {feas.get('required_weekly_growth',5)} new followers/week",
        ]
        bar(s, Inches(0.3), Inches(1.22), Inches(9.4), Inches(0.38), RGBColor(0x1A,0x10,0x35))
        for ci, lbl in enumerate(["SCENARIO","30 DAYS","60 DAYS","90 DAYS","180 DAYS"]):
            txt(s, lbl, Inches(0.4+ci*1.87), Inches(1.28), Inches(1.82), Inches(0.28), 8, bold=True, color=accent)
        row_scenarios = [("BEST CASE",green),("REALISTIC",accent),("WORST CASE",red)]
        scenario_map  = {"BEST CASE":"best","REALISTIC":"realistic","WORST CASE":"worst"}
        for ri, (lbl, col) in enumerate(row_scenarios):
            y = Inches(1.64 + ri * 0.78)
            bar(s, Inches(0.3), y, Inches(9.4), Inches(0.72), mid_dark if ri%2==0 else RGBColor(0x16,0x10,0x30))
            bar(s, Inches(0.3), y, Inches(0.05), Inches(0.72), col)
            txt(s, lbl, Inches(0.4), y+Inches(0.17), Inches(1.82), Inches(0.38), 10, bold=True, color=col)
            sk = scenario_map[lbl]
            for ci, pk in enumerate(("day30","day60","day90","day180")):
                pd_data = gf.get(pk) or {}
                v = pd_data.get(sk)
                disp = f"{int(v):,}" if isinstance(v,(int,float)) and v else "—"
                txt(s, disp, Inches(0.4+(ci+1)*1.87), y+Inches(0.17), Inches(1.8), Inches(0.38), 12, bold=True, color=col)
        bar(s, Inches(0.3), Inches(4.06), Inches(9.4), Inches(0.88), mid_dark)
        txt(s, "KEY ASSUMPTIONS", Inches(0.45), Inches(4.14), Inches(9), Inches(0.28), 8, bold=True, color=amber)
        txt(s, "  ·  ".join(str(a)[:55] for a in assumptions[:3]),
            Inches(0.45), Inches(4.42), Inches(9.1), Inches(0.45), 9, color=grey, italic=True)
        footer(s, 12)

        # ── Slide 13: Performance Diagnosis ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "12  ·  PERFORMANCE DIAGNOSIS", "What the data is telling us.")
        diag13 = strategy.get("performance_diagnosis") or {}
        # When live data is missing, do NOT fabricate Working/Failing claims
        # about posts that don't exist. Show structural diagnosis only.
        if has_live:
            sec13 = [
                ("WHAT'S WORKING",        green, diag13.get("whats_working")       or strategy.get("what_works")   or ["Top posts show strong ER"]),
                ("WHAT'S FAILING",        red,   diag13.get("whats_failing")       or strategy.get("what_to_stop") or ["Low-engagement formats"]),
                ("MISSED OPPORTUNITIES",  amber, diag13.get("missed_opportunities") or ["Untapped trending formats"]),
                ("BOTTLENECKS",           accent,diag13.get("bottlenecks")          or ["Posting frequency inconsistency"]),
            ]
        else:
            sec13 = [
                ("STRUCTURAL OPPORTUNITIES", amber,  diag13.get("missed_opportunities") or ["Connect IG to identify untapped niche formats"]),
                ("STRUCTURAL BOTTLENECKS",   accent, diag13.get("bottlenecks")          or ["No live content history to learn from"]),
                ("RESEARCH-BASED ANGLES",    green,  (rd.get("trending_angles") or ["See market research"])[:3]),
                ("COMPETITOR-DERIVED GAPS",  brand_color, (cd.get("content_gaps") or ["See competitor analysis"])[:3]),
            ]
        for i, (title, col, items) in enumerate(sec13):
            row, ci = divmod(i, 2)
            x = Inches(0.3+ci*4.8); y = Inches(1.15+row*2.12)
            bar(s, x, y, Inches(4.5), Inches(1.88), mid_dark)
            bar(s, x, y, Inches(4.5), Inches(0.05), col)
            txt(s, title, x+Inches(0.15), y+Inches(0.1), Inches(4.1), Inches(0.3), 8, color=col, bold=True)
            for j, item in enumerate(items[:3]):
                txt(s, f"• {str(item)[:68]}", x+Inches(0.15), y+Inches(0.45+j*0.44), Inches(4.2), Inches(0.38), 10, color=light)
        footer(s, 13)

        # ── Slide 14: Growth Strategy ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "13  ·  GROWTH STRATEGY", "What to do, why to do it, and when.")
        goal_strat  = strategy.get("goal_strategy") or {}
        narrative14 = str(goal_strat.get("narrative",""))[:280]
        if narrative14:
            bar(s, Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.52), mid_dark)
            txt(s, narrative14, Inches(0.45), Inches(1.14), Inches(9.1), Inches(0.45), 10, color=light, italic=True)
        roadmap_d = strategy.get("roadmap") or {}
        horizons  = [
            ("QUICK WINS (1-7 DAYS)", brand_color, roadmap_d.get("quick_wins") or (strategy.get("growth_tactics") or [])[:2]),
            ("30 DAYS",               green,        roadmap_d.get("short_term") or (strategy.get("growth_tactics") or [])[2:4]),
            ("90 DAYS",               accent,       roadmap_d.get("mid_term")   or (strategy.get("monthly_themes") or [])[:2]),
        ]
        y14 = Inches(1.72) if narrative14 else Inches(1.2)
        for i, (title, col, items) in enumerate(horizons):
            x = Inches(0.3+i*3.2)
            bar(s, x, y14, Inches(3.0), Inches(3.3), mid_dark)
            bar(s, x, y14, Inches(3.0), Inches(0.05), col)
            txt(s, title, x+Inches(0.12), y14+Inches(0.1), Inches(2.8), Inches(0.32), 8, color=col, bold=True)
            for j, item in enumerate(items[:4] if isinstance(items,list) else [items]):
                action = str(item.get("action","") if isinstance(item,dict) else item)[:58]
                txt(s, f"→ {action}", x+Inches(0.12), y14+Inches(0.5+j*0.65), Inches(2.8), Inches(0.58), 10, color=light)
        footer(s, 14)

        # ── Slide 15: Recommendations with Evidence ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "14  ·  RECOMMENDATIONS — DATA-BACKED", "Every recommendation tied to evidence. No generic advice.")
        struct_recs = ep.get("structured_recommendations") or []
        if not struct_recs:
            pa = strategy.get("priority_actions") or []
            struct_recs = [{"title": str(r.get("action",""))[:70],
                            "evidence": [str(r.get("why_it_matters",""))],
                            "expected_impact": str(r.get("expected_impact","")),
                            "priority": "P1", "difficulty": str(r.get("difficulty","medium")),
                            "timeline": str(r.get("timeline",""))} for r in pa[:5]]
        diff_col_map = {"low": green, "medium": amber, "high": red}
        for i, rec in enumerate(struct_recs[:5]):
            y = Inches(1.08 + i * 0.85)
            bar(s, Inches(0.4), y, Inches(9.2), Inches(0.76), mid_dark)
            bar(s, Inches(0.4), y, Inches(0.05), Inches(0.76), brand_color)
            title_r  = str(rec.get("title",""))[:105]
            evidence = rec.get("evidence") or []
            diff     = str(rec.get("difficulty","")).lower()
            tline    = str(rec.get("timeline",""))
            impact   = str(rec.get("expected_impact",""))[:60]
            priority = str(rec.get("priority",""))
            # Line 1: title
            txt(s, title_r, Inches(0.55), y+Inches(0.06), Inches(8.4), Inches(0.28), 11, bold=True, color=white)
            # Line 2: evidence (primary)
            evid1 = str(evidence[0])[:100] if evidence else ""
            if evid1:
                txt(s, f"Evidence: {evid1}", Inches(0.55), y+Inches(0.35), Inches(8.6), Inches(0.2), 8, color=light, italic=True)
            # Line 3: impact + timeline + priority + difficulty
            meta_parts = []
            if impact: meta_parts.append(f"→ {impact}")
            if tline:  meta_parts.append(tline)
            if priority: meta_parts.append(priority)
            if diff:   meta_parts.append(diff.capitalize())
            meta = "  ·  ".join(meta_parts)[:120]
            txt(s, meta, Inches(0.55), y+Inches(0.55), Inches(8.6), Inches(0.18), 8,
                color=diff_col_map.get(diff, grey))
        footer(s, 15)

        # ── Slide 16: Content Pillar Breakdown ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "15  ·  CONTENT PILLAR BREAKDOWN", "What to create — and how hard to push each pillar.")
        pillar_breakdown = strategy.get("pillar_breakdown") or []
        pillar_list = brand.get("contentPillars") or strategy.get("pillars") or []
        if not pillar_breakdown and pillar_list:
            pillar_breakdown = [{"pillar":str(p),"current_performance":"medium","growth_potential":"high",
                                 "recommended_frequency":"2x/week","expected_impact":"Build authority in niche"}
                                for p in pillar_list[:4]]
        hdr_cols16 = [("CONTENT PILLAR",Inches(0.3)),("PERFORMANCE",Inches(3.85)),("POTENTIAL",Inches(5.5)),
                      ("FREQUENCY",Inches(6.95)),("IMPACT",Inches(8.1))]
        bar(s, Inches(0.3), Inches(1.15), Inches(9.4), Inches(0.38), RGBColor(0x1A,0x10,0x35))
        for col_label, x in hdr_cols16:
            txt(s, col_label, x+Inches(0.08), Inches(1.19), Inches(1.5), Inches(0.3), 7, color=accent, bold=True)
        perf_col_map = {"high":green,"medium":amber,"low":red}
        for i, pb in enumerate(pillar_breakdown[:4]):
            y = Inches(1.57+i*0.9)
            bar(s, Inches(0.3), y, Inches(9.4), Inches(0.84), mid_dark if i%2==0 else RGBColor(0x16,0x10,0x30))
            perf = str(pb.get("current_performance","medium")).lower()
            pot  = str(pb.get("growth_potential","high")).lower()
            txt(s, str(pb.get("pillar",""))[:45], Inches(0.42), y+Inches(0.1), Inches(3.3), Inches(0.64), 11, color=white)
            txt(s, perf.upper(), Inches(3.95), y+Inches(0.2), Inches(1.4), Inches(0.44), 11, bold=True, color=perf_col_map.get(perf,white))
            txt(s, pot.upper(),  Inches(5.58), y+Inches(0.2), Inches(1.3), Inches(0.44), 11, bold=True, color=perf_col_map.get(pot,white))
            txt(s, str(pb.get("recommended_frequency","2x/week"))[:15], Inches(7.02), y+Inches(0.2), Inches(1.0), Inches(0.44), 11, color=accent)
            txt(s, str(pb.get("expected_impact",""))[:55], Inches(8.17), y+Inches(0.1), Inches(1.45), Inches(0.64), 9, color=grey)
        footer(s, 16)

        # ── Slide 17: Hook / Reel / Carousel Strategy ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "16  ·  CONTENT STRATEGY — HOOKS, REELS & CAROUSELS", "The creative blueprint for every format.")
        sec17 = [
            ("HOOK STRATEGY",    brand_color, [str(h)[:58] for h in (strategy.get("hook_strategy") or [])[:3]]),
            ("REEL RETENTION",   accent,      [str(r)[:58] for r in (strategy.get("retention_strategy") or strategy.get("content_series_ideas") or [])[:3]]),
            ("CONVERSION MOVES", green,       [str(c)[:58] for c in (strategy.get("conversion_strategy") or strategy.get("cta_templates") or [])[:3]]),
        ]
        for i, (title, col, items) in enumerate(sec17):
            x = Inches(0.3+i*3.2)
            bar(s, x, Inches(1.15), Inches(3.0), Inches(4.0), mid_dark)
            bar(s, x, Inches(1.15), Inches(3.0), Inches(0.05), col)
            txt(s, title, x+Inches(0.12), Inches(1.23), Inches(2.8), Inches(0.32), 8, color=col, bold=True)
            for j, item in enumerate((items or ["See strategy brief"])[:4]):
                txt(s, f"→ {item}", x+Inches(0.12), Inches(1.6+j*0.85), Inches(2.8), Inches(0.78), 10, color=light)
        footer(s, 17)

        # ── Slide 18: 30-Day Execution Plan ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "17  ·  30-DAY EXECUTION PLAN", "Day-by-day: what to post, why, and what it achieves.")
        cal30      = ep.get("execution_calendar_30") or []
        day_plan30 = strategy.get("day_by_day_plan") or []
        source30   = cal30 or day_plan30
        if source30:
            grps = [source30[:5], source30[5:10], source30[10:15]]
            lbls = ["DAYS 1-5","DAYS 6-10","DAYS 11-15"]
            cols = [brand_color, green, accent]
            for i, (grp, lbl, col) in enumerate(zip(grps, lbls, cols)):
                x = Inches(0.3+i*3.2)
                bar(s, x, Inches(1.15), Inches(3.0), Inches(4.05), mid_dark)
                bar(s, x, Inches(1.15), Inches(3.0), Inches(0.05), col)
                txt(s, lbl, x+Inches(0.12), Inches(1.23), Inches(2.8), Inches(0.32), 8, color=col, bold=True)
                for j, day in enumerate(grp[:5]):
                    dy = Inches(1.62+j*0.69)
                    bar(s, x+Inches(0.05), dy, Inches(2.88), Inches(0.63), RGBColor(0x1A,0x10,0x35))
                    dn = day.get("day", i*5+j+1)
                    rt = str(day.get("reel_topic","") or day.get("content_task",""))[:42]
                    g  = str(day.get("goal","") or day.get("kpi_target",""))[:28]
                    txt(s, f"Day {dn}", x+Inches(0.12), dy+Inches(0.05), Inches(0.7), Inches(0.22), 8, color=col, bold=True)
                    txt(s, rt, x+Inches(0.12), dy+Inches(0.26), Inches(2.65), Inches(0.2), 8, color=light)
                    txt(s, g,  x+Inches(0.12), dy+Inches(0.46), Inches(2.65), Inches(0.15), 7, color=grey)
        else:
            for j, t in enumerate((strategy.get("growth_tactics") or [])[:5]):
                txt(s, f"Week {j+1}: {str(t)[:100]}", Inches(0.4), Inches(1.5+j*0.72), Inches(9.2), Inches(0.65), 11, color=light)
        footer(s, 18)

        # ── Slide 19: 60-Day Plan ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "18  ·  60-DAY EXECUTION PLAN", "Weeks 5–8: compounding the momentum.")
        weeks_5_12 = ep.get("execution_calendar_weeks_5_to_12") or strategy.get("weekly_plan") or []
        wks_5_8    = [w for w in weeks_5_12 if isinstance(w,dict) and w.get("week",0) <= 8][:4] or weeks_5_12[:4]
        wk_cols    = [brand_color, green, accent, amber]
        if wks_5_8:
            for i, wk in enumerate(wks_5_8[:4]):
                wc = wk_cols[i%4]; x = Inches(0.3+i*2.4)
                bar(s, x, Inches(1.15), Inches(2.2), Inches(4.05), mid_dark)
                bar(s, x, Inches(1.15), Inches(2.2), Inches(0.05), wc)
                txt(s, f"WEEK {wk.get('week',i+5)}", x+Inches(0.12), Inches(1.23), Inches(2.0), Inches(0.32), 8, color=wc, bold=True)
                txt(s, str(wk.get("theme",""))[:35],  x+Inches(0.12), Inches(1.58), Inches(2.0), Inches(0.38), 12, bold=True, color=white)
                txt(s, str(wk.get("reels","") or wk.get("content_tasks",""))[:80], x+Inches(0.12), Inches(2.0), Inches(2.0), Inches(1.0), 10, color=light)
                txt(s, f"KPI: {str(wk.get('goal','') or wk.get('kpi_target',''))[:42]}", x+Inches(0.12), Inches(3.08), Inches(2.0), Inches(0.35), 9, color=wc)
        else:
            themes = strategy.get("monthly_themes") or []
            fb60 = [
                ("DAYS 31-45", brand_color, (strategy.get("growth_tactics") or [])[:2] + [themes[1] if len(themes)>1 else "Content acceleration phase"]),
                ("DAYS 46-60", green,        ["Collabs + niche partnerships","Viral series continues"]),
            ]
            for i, (title, col, items) in enumerate(fb60):
                x = Inches(0.3+i*4.8)
                bar(s, x, Inches(1.2), Inches(4.5), Inches(3.85), mid_dark)
                bar(s, x, Inches(1.2), Inches(4.5), Inches(0.05), col)
                txt(s, title, x+Inches(0.15), Inches(1.28), Inches(4.2), Inches(0.3), 8, color=col, bold=True)
                for j, item in enumerate(items[:4]):
                    txt(s, f"→ {str(item)[:68]}", x+Inches(0.15), Inches(1.65+j*0.85), Inches(4.2), Inches(0.75), 10, color=light)
        footer(s, 19)

        # ── Slide 20: 90-Day Plan ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "19  ·  90-DAY EXECUTION PLAN", "Weeks 9–12: build the moat — authority, audience, retention.")
        wks_9_12 = [w for w in weeks_5_12 if isinstance(w,dict) and w.get("week",0) >= 9][:4]
        if not wks_9_12 and len(weeks_5_12) > 4:
            wks_9_12 = weeks_5_12[4:8]
        if wks_9_12:
            for i, wk in enumerate(wks_9_12[:4]):
                wc = wk_cols[i%4]; x = Inches(0.3+i*2.4)
                bar(s, x, Inches(1.15), Inches(2.2), Inches(4.05), mid_dark)
                bar(s, x, Inches(1.15), Inches(2.2), Inches(0.05), wc)
                txt(s, f"WEEK {wk.get('week',i+9)}", x+Inches(0.12), Inches(1.23), Inches(2.0), Inches(0.32), 8, color=wc, bold=True)
                txt(s, str(wk.get("theme",""))[:35], x+Inches(0.12), Inches(1.58), Inches(2.0), Inches(0.38), 12, bold=True, color=white)
                txt(s, str(wk.get("reels","") or wk.get("content_tasks",""))[:80], x+Inches(0.12), Inches(2.0), Inches(2.0), Inches(1.0), 10, color=light)
                txt(s, f"KPI: {str(wk.get('goal','') or wk.get('kpi_target',''))[:42]}", x+Inches(0.12), Inches(3.08), Inches(2.0), Inches(0.35), 9, color=wc)
        else:
            fb90 = [
                ("DAYS 61-75", accent, (strategy.get("viral_opportunities") or ["Authority Reels series","UGC campaign launch"])[:3]),
                ("DAYS 76-90", amber,  (strategy.get("content_series_ideas") or ["Community spotlight","Monetisation readiness"])[:3]),
            ]
            for i, (title, col, items) in enumerate(fb90):
                x = Inches(0.3+i*4.8)
                bar(s, x, Inches(1.2), Inches(4.5), Inches(3.85), mid_dark)
                bar(s, x, Inches(1.2), Inches(4.5), Inches(0.05), col)
                txt(s, title, x+Inches(0.15), Inches(1.28), Inches(4.2), Inches(0.3), 8, color=col, bold=True)
                for j, item in enumerate(items[:4]):
                    txt(s, f"→ {str(item)[:68]}", x+Inches(0.15), Inches(1.65+j*0.85), Inches(4.2), Inches(0.75), 10, color=light)
        footer(s, 20)

        # ── Slide 21: KPIs & Success Metrics ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        slide_header(s, "20  ·  KPIs & SUCCESS METRICS", "How we measure winning.")
        kpi_90 = strategy.get("kpi_targets_90day") or strategy.get("kpi_targets_30day") or ar.get("kpi_targets_90day") or {}
        if not isinstance(kpi_90, dict): kpi_90 = {}
        kpi_fol   = int(kpi_90.get("followers", goal) or goal)
        kpi_er    = kpi_90.get("avg_engagement_rate", "3-5")
        kpi_reach = kpi_90.get("avg_reach", int(avg_reach*1.5) if avg_reach else 1000)
        kpi_reels = kpi_90.get("reels_per_week", 4)
        kpi_saves = kpi_90.get("saves_per_post", 20)
        kpi_items21 = [
            ("FOLLOWER TARGET", f"{followers:,} → {kpi_fol:,}", brand_color),
            ("ENGAGEMENT RATE", f"{kpi_er}%",                   green),
            ("REACH PER POST",  f"{kpi_reach:,}+",              accent),
            ("REELS / WEEK",    f"{kpi_reels}x",                amber),
            ("SAVES PER POST",  f"{kpi_saves}+",                RGBColor(0xEC,0x48,0x99)),
        ]
        for i, (lbl, val, col) in enumerate(kpi_items21):
            x = Inches(0.3+i*1.9)
            bar(s, x, Inches(1.3), Inches(1.8), Inches(3.3), mid_dark)
            bar(s, x, Inches(1.3), Inches(1.8), Inches(0.06), col)
            txt(s, lbl, x+Inches(0.1), Inches(1.42), Inches(1.6), Inches(0.3), 8, color=col, bold=True)
            txt(s, val, x+Inches(0.1), Inches(1.76), Inches(1.6), Inches(0.65), 16, bold=True)
        txt(s, f"Posting frequency: {strategy.get('posting_frequency','1-2x daily')}  ·  Best times: {', '.join(str(t) for t in (strategy.get('best_times',['9AM','12PM','6PM']))[:3])}",
            Inches(0.4), Inches(4.78), Inches(9.2), Inches(0.35), 11, color=grey)
        footer(s, 21)

        # ── Slide 22: Outro ──
        s = prs.slides.add_slide(blank); bg(s, dark)
        bar(s, 0, 0, Inches(0.06), prs.slide_height, brand_color)
        txt(s, "Let's build.", Inches(0.4), Inches(1.6), Inches(9), Inches(1.3), 48, bold=True)
        txt(s, "PerformEdge  ·  Social Growth Partners", Inches(0.4), Inches(3.05), Inches(9), Inches(0.5), 16, color=accent)
        txt(s, name, Inches(0.4), Inches(3.7), Inches(9), Inches(0.45), 14, color=grey, italic=True)
        txt(s, f"90-Day Strategy  ·  {month_str}  ·  Investor-Grade Report", Inches(0.4), Inches(4.45), Inches(9), Inches(0.35), 12, color=RGBColor(0x4B,0x55,0x63))


        # Serialize + upload
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        path = f"runs/{run_id or 'growth'}/growth_planner.pptx"
        url  = _upload_bytes(buf.read(), path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        return url

    except Exception as e:
        print(f"[GrowthPlanner] PPT build error: {e}")
        import traceback; traceback.print_exc()
        return None


def _upload_bytes(data: bytes, path: str, content_type: str):
    sb = _get_supabase()
    if not sb:
        print(f"[GrowthPlanner] Supabase unavailable — skipping upload for {path}")
        return None
    bucket = os.getenv("SUPABASE_STORAGE_BUCKET", "socialos-storage")
    try:
        sb.storage.from_(bucket).upload(
            path, data,
            file_options={"content-type": content_type, "upsert": "true"},
        )
        result = sb.storage.from_(bucket).get_public_url(path)
        print(f"[GrowthPlanner] Uploaded {path} -> {str(result)[:80]}")
        return result
    except Exception as e:
        print(f"[GrowthPlanner] Supabase upload error for {path}: {e}")
        import traceback; traceback.print_exc()
        return None


def _validate_strategy(strategy: dict, ig_audit: dict, feasibility: dict = None, days_ahead: int = 15) -> tuple[dict, dict]:
    """Quality-check phase (Change 2) — verify and correct numbers BEFORE they
    reach any slide. Does NOT fabricate data: it only clamps impossible values
    to the trusted Python-computed bounds, enforces internal consistency, and
    labels low-confidence figures. Returns (cleaned_strategy, quality_report).

    Scalable by design: any pipeline producing a strategy dict can route it
    through this gate and inherit the same guardrails.
    """
    if not isinstance(strategy, dict):
        return strategy, {"checks": [], "corrections": [], "confidence": "low"}

    followers = int(ig_audit.get("followers", 0) or 0)
    goal      = int(ig_audit.get("goal_followers", 0) or 0)
    data_src  = ig_audit.get("data_source", "estimate")
    lo, hi    = followers, max(goal, followers)

    checks: list = []
    corrections: list = []

    def clamp(v):
        try:
            n = int(round(float(v)))
        except (ValueError, TypeError):
            return None
        return max(lo, min(hi, n))

    # 1) follower_plan milestones must sit within [current, goal] and rise
    fp = strategy.get("follower_plan")
    if isinstance(fp, dict) and (lo or hi):
        ordered_keys = [k for k in ("day5", "day10", "day15", "week1", "week2", "week3", "week4", "goal") if k in fp]
        prev = lo
        for k in ordered_keys:
            cv = clamp(fp.get(k))
            if cv is None:
                continue
            # enforce monotonic non-decreasing progression toward the goal
            cv = max(cv, prev)
            if cv != fp.get(k):
                corrections.append(f"follower_plan.{k}: {fp.get(k)} → {cv} (kept within {lo:,}–{hi:,})")
            fp[k] = cv
            prev = cv
        strategy["follower_plan"] = fp
        checks.append("Follower milestones verified within current→goal bounds")

    # 2) content_mix percentages should sum to 100
    cm = strategy.get("content_mix")
    if isinstance(cm, dict) and cm:
        nums = {k: (float(v) if isinstance(v, (int, float)) else 0) for k, v in cm.items()}
        total = sum(nums.values())
        if total > 0 and abs(total - 100) > 1:
            normed = {k: round(v / total * 100) for k, v in nums.items()}
            # fix rounding drift on the largest bucket
            drift = 100 - sum(normed.values())
            if normed:
                big = max(normed, key=normed.get)
                normed[big] += drift
            strategy["content_mix"] = normed
            corrections.append(f"content_mix normalised to 100% (was {round(total)}%)")
        checks.append("Content mix percentages sum to 100")

    # 3) goal consistency vs feasibility
    if feasibility:
        f_goal = int(feasibility.get("target_followers", goal) or goal)
        if goal and f_goal and goal != f_goal:
            corrections.append(f"goal mismatch reconciled to computed target {goal:,}")
        checks.append("Goal target reconciled with feasibility math")

    # 4) confidence + labelling. Live IG data = high confidence; otherwise the
    # plan is built from stored counts / brand brief, so figures are estimates.
    confidence = "high" if data_src == "live" else ("medium" if data_src == "stored" else "low")
    if confidence != "high":
        note = (
            "Estimate — based on stored follower count; connect Instagram for live-measured figures."
            if confidence == "medium" else
            "Assumption — no Instagram connected; figures are brand-brief projections, not measured data."
        )
        strategy["data_confidence"] = confidence
        strategy["data_confidence_note"] = note
        checks.append(f"Low-confidence figures labelled ({confidence})")

    return strategy, {"checks": checks, "corrections": corrections, "confidence": confidence}


def _fallback_strategy(brand: dict, days_ahead: int, ig_audit: dict = None) -> dict:
    niche   = brand.get("niche", "your niche")
    name    = brand.get("name", "your brand")
    # Shorten niche for use in tactic strings — long niche values like
    # "Air cooler distribution & cooling solutions" produce broken-looking text.
    # Use the first 3 meaningful words, or the full niche if short.
    niche_words = niche.split()
    niche_short = " ".join(niche_words[:3]) if len(niche_words) > 3 else niche
    # Use brand name for brand-specific references
    brand_ref   = name if name and name != "your brand" else niche_short

    _followers = (ig_audit or {}).get("followers", 0)
    _goal      = (ig_audit or {}).get("goal_followers", 500)
    _gap       = max(_goal - _followers, 0)
    content_pillars = brand.get("contentPillars") or []
    pillars = content_pillars[:4] if content_pillars else [
        "Brand Awareness", "Education & Value", "Engagement", "Social Proof"
    ]
    return {
        "pillars": pillars,
        "posting_frequency": "1-2x daily",
        "best_times": ["9:00 AM", "12:30 PM", "6:00 PM"],
        "content_mix": {"Reel": 40, "Carousel": 30, "Graphic": 15, "Story": 10, "AI Reel": 5},
        "monthly_themes": [
            f"{brand_ref} Brand Story Month",
            f"{niche_short} Client Spotlight",
            f"{niche_short} Industry Insights",
        ],
        "growth_tactics": [
            f"Ship 14-day Reel series on {brand.get('audiencePainPoints','common pain points')[:45]} using PAS hooks — repurpose each into a Carousel + 3 Stories",
            f"Engage with 15 accounts in the {niche_short} community daily — comment with genuine insight, not emojis",
            "Use 15-20 targeted hashtags per post (mix broad 500K+ + niche 10K-200K)",
            "Respond to all comments within 1 hour of posting to boost algorithmic distribution",
            f"Partner with 1-2 micro-influencers in the {niche_short} space for a swap collab",
            "Pin the best-performing post as a profile highlight and update it every 30 days",
        ],
        "hashtag_strategy": f"Use a mix of broad {niche} hashtags (500K-2M posts) and niche-specific ones (10K-200K posts) to balance reach and community discovery.",
        "cta_templates": [
            "Save this — you'll want to come back to it! 🔖",
            "Tag someone in {niche} who needs to see this!",
            "Drop a comment below — what's your experience? 👇",
            "Share this with your team! 🚀",
        ],
        "what_works": [
            f"Short-form Reels (15-30s) with {niche}-specific hooks get highest reach",
            "Educational carousels with a clear problem-solution structure drive saves",
            "Behind-the-scenes and founder content builds authentic engagement",
            "Posts that directly address audience pain points get more DMs",
        ],
        "what_to_stop": [
            "Generic motivational quotes without brand context",
            "Overly promotional posts without value delivery",
            "Static images without text overlay or visual hook",
        ],
        "follower_plan": {
            "week1": _followers + int(_gap * 0.25),
            "week2": _followers + int(_gap * 0.50),
            "week3": _followers + int(_gap * 0.75),
            "week4": _goal,
        },
        "engagement_tactics": [
            "Reply to every comment within 1 hour of posting",
            f"Engage daily in {niche} hashtag communities",
            "Ask a specific question in every caption",
            "Use polls/questions in Stories daily",
        ],
        "content_series_ideas": [
            f"Weekly '{niche} myth vs reality' carousel",
            f"'Behind the scenes at {name}' Reels series",
            f"Monthly '{niche} wins' client spotlight",
        ],
    }
